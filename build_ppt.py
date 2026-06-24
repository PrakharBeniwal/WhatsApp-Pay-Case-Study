# -*- coding: utf-8 -*-
"""
WhatsApp Pay India - Product Strategy Case Study
12-slide, visual-first deck built entirely with python-pptx (no external images).
Usage: python build_ppt.py
Output: WhatsApp_Pay_Case_Study.pptx (same folder)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# ---------- palette ----------
def C(h): return RGBColor.from_string(h)
PRIMARY   = C('075E54')
SECONDARY = C('128C7E')
ACCENT    = C('25D366')
TEAL      = C('34B7F1')
DARK      = C('0B141A')
LIGHT     = C('F0F2F5')
WHITE     = C('FFFFFF')
BORDER    = C('E0E0E0')
RED       = C('E63946')
AMBER     = C('F4A261')
GREY      = C('8D99AE')
INK       = C('1A1A1A')
CARD_DK   = C('10212B')
AMBER_TINT= C('FCEEDD')
RED_TINT  = C('FBE3E5')
GREEN_TINT= C('E8F7F0')

# pyramid greens (dark -> light, top -> bottom)
PYR = [C('075E54'), C('0F7266'), C('128C7E'), C('1AA890'), C('25C49A')]

EMU = 914400

# ---------- low-level helpers ----------
def R(t, s, b, c, i=False):
    return (t, s, b, c, i)

def P(runs, align=None, sa=None, sb=None, ln=None):
    d = {'runs': runs}
    if align is not None: d['align'] = align
    if sa is not None: d['sa'] = sa
    if sb is not None: d['sb'] = sb
    if ln is not None: d['ln'] = ln
    return d

def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align)
        if 'sa' in para: p.space_after = Pt(para['sa'])
        if 'sb' in para: p.space_before = Pt(para['sb'])
        if 'ln' in para: p.line_spacing = para['ln']
        for run in para['runs']:
            r = p.add_run(); r.text = run[0]
            r.font.size = Pt(run[1]); r.font.bold = run[2]
            r.font.color.rgb = run[3]; r.font.name = 'Calibri'
            if len(run) > 4 and run[4]: r.font.italic = True
    return box

def card(slide, x, y, w, h, fill=WHITE, line=BORDER, line_w=1.0, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    return shp

def chip(slide, x, y, w, h, txt, fill, txt_color=WHITE, size=11, bold=True,
         align=PP_ALIGN.CENTER, radius=0.5, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = txt_color; r.font.name = 'Calibri'
    return shp

def icon(slide, x, y, d, glyph, fill, gcolor=WHITE, size=16, line=None):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None: o.line.fill.background()
    else: o.line.color.rgb = line; o.line.width = Pt(1.5)
    o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = gcolor; r.font.name = 'Calibri'
    return o

def arrow_line(slide, x1, y1, x2, y2, color, weight=2.5, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    if dash:
        ln.append(parse_xml('<a:prstDash %s val="%s"/>' % (nsdecls('a'), dash)))
    ln.append(parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>' % nsdecls('a')))
    conn.shadow.inherit = False
    return conn

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def add_header(slide, title, kicker):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background(); band.shadow.inherit = False
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0), Inches(13.333), Inches(0.06))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background(); rule.shadow.inherit = False
    tb(slide, 0.5, 0.1, 9.7, 0.82, [P([R(title, 25, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
    chip(slide, 10.85, 0.3, 2.0, 0.42, kicker, SECONDARY, WHITE, 12, True, radius=0.5)

def add_footer(slide, n, dark=False):
    col = GREY
    tb(slide, 0.5, 7.08, 8.0, 0.34,
       [P([R('WhatsApp Pay India  —  Product Strategy Case Study', 9, False, col)])],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, 10.8, 7.08, 2.03, 0.34,
       [P([R('%d / 12' % n, 10, True, col)], align=PP_ALIGN.RIGHT)],
       anchor=MSO_ANCHOR.MIDDLE)

# ---------- chart helpers ----------
def _style_chart(chart, font_sz=11):
    try:
        chart.font.name = 'Calibri'; chart.font.size = Pt(font_sz)
    except Exception: pass
    chart.has_title = False

def color_points(series, colors):
    for idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[idx % len(colors)]
        try: pt.format.line.fill.background()
        except Exception: pass

def native_bar(slide, x, y, w, h, cats, vals, colors, maxval=None,
               number_format='0.0', cat_sz=10, lbl_sz=11):
    cats = list(cats)[::-1]; vals = list(vals)[::-1]; colors = list(colors)[::-1]
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('S', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart; _style_chart(chart, lbl_sz)
    chart.has_legend = False
    plot = chart.plots[0]; plot.gap_width = 55
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format; dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(lbl_sz); dl.font.bold = True; dl.font.name = 'Calibri'; dl.font.color.rgb = INK
    color_points(plot.series[0], colors)
    va = chart.value_axis
    va.minimum_scale = 0
    if maxval: va.maximum_scale = maxval
    va.visible = False
    try: va.has_major_gridlines = False
    except Exception: pass
    ca = chart.category_axis
    try:
        ca.tick_labels.font.size = Pt(cat_sz); ca.tick_labels.font.name = 'Calibri'
        ca.tick_labels.font.color.rgb = INK
    except Exception: pass
    try: ca.format.line.color.rgb = BORDER
    except Exception: pass
    return chart

def native_pie(slide, x, y, w, h, cats, vals, colors):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('S', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart; _style_chart(chart, 11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11); chart.legend.font.name = 'Calibri'
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0"%"'; dl.number_format_is_linked = False
    dl.show_value = True; dl.show_percentage = False; dl.show_category_name = False
    dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = 'Calibri'; dl.font.color.rgb = WHITE
    color_points(plot.series[0], colors)
    return chart

def native_col(slide, x, y, w, h, cats, vals, colors, labels, maxval, cat_sz=11):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('S', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart; _style_chart(chart, 12)
    chart.has_legend = False
    plot = chart.plots[0]; plot.gap_width = 70
    plot.has_data_labels = True
    color_points(plot.series[0], colors)
    for idx, pt in enumerate(plot.series[0].points):
        dl = pt.data_label
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        tf = dl.text_frame; tf.text = labels[idx]
        r = tf.paragraphs[0].runs[0]
        r.font.size = Pt(14); r.font.bold = True; r.font.name = 'Calibri'; r.font.color.rgb = PRIMARY
    va = chart.value_axis
    va.minimum_scale = 0; va.maximum_scale = maxval; va.visible = False
    try: va.has_major_gridlines = False
    except Exception: pass
    ca = chart.category_axis
    try:
        ca.tick_labels.font.size = Pt(cat_sz); ca.tick_labels.font.bold = True
        ca.tick_labels.font.name = 'Calibri'; ca.tick_labels.font.color.rgb = INK
    except Exception: pass
    return chart

# =====================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------------- SLIDE 1 : TITLE / HERO ----------------
    s = prs.slides.add_slide(blank); bg(s, DARK)
    # faint rupee watermark
    tb(s, 4.0, 0.55, 5.33, 5.6, [P([R('₹', 360, True, PRIMARY)], align=PP_ALIGN.CENTER)],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 0.8, 1.85, 11.73, 1.9, [
        P([R('500 Million Users.  ₹0 Revenue.', 40, True, WHITE)], align=PP_ALIGN.CENTER, sa=4),
        P([R('Why WhatsApp Pay Lost India.', 40, True, ACCENT)], align=PP_ALIGN.CENTER),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.17), Inches(3.78), Inches(3.0), Inches(0.04))
    r2.fill.solid(); r2.fill.fore_color.rgb = ACCENT; r2.line.fill.background(); r2.shadow.inherit = False
    tb(s, 1.5, 3.95, 10.33, 0.7, [P([R(
        'A Product Strategy Case Study  —  Competitive Teardown  →  Root-Cause Diagnosis  →  Rejuvenation Strategy',
        15, False, WHITE)], align=PP_ALIGN.CENTER)], align=PP_ALIGN.CENTER)
    stats = ['3 PHASES', '84,000+ WORDS', '51 SOLUTIONS', 'RICE-PRIORITIZED']
    cw, gap = 2.6, 0.2
    x0 = (13.333 - (4 * cw + 3 * gap)) / 2
    for i, st in enumerate(stats):
        chip(s, x0 + i * (cw + gap), 4.85, cw, 0.52, st, SECONDARY, WHITE, 13, True, radius=0.5)
    tb(s, 0.6, 6.85, 6.0, 0.4, [P([R('[Your Name]  ·  PM Portfolio', 12, False, GREY)])],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 6.7, 6.85, 6.03, 0.4, [P([R('fahadahmed982004@gmail.com', 12, False, GREY)], align=PP_ALIGN.RIGHT)],
       anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 2 : THE PARADOX ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'The Paradox', 'THE SETUP'); add_footer(s, 2)
    L = [('500M+ users', 'most-used app in India'),
         ('Zero CAC', 'pre-installed audience'),
         ('$19B', "Meta's backing"),
         ('Chat-native', 'no app-switch friction'),
         ('15M', 'WhatsApp Business accounts')]
    Rr = [('<1%', 'UPI market share'),
          ('Rank #4', 'behind PhonePe, GPay, Paytm'),
          ('<10M', 'active P2M payers (of ~100M)'),
          ('₹0', 'revenue under zero-MDR'),
          ('Near-zero', 'merchant acceptance')]
    # left card
    card(s, 0.5, 1.35, 5.7, 4.65, WHITE, BORDER)
    card(s, 0.5, 1.35, 5.7, 0.55, PRIMARY, None)
    tb(s, 0.5, 1.35, 5.7, 0.55, [P([R('WHAT WHATSAPP HAD', 14, True, WHITE)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    # right card
    card(s, 7.13, 1.35, 5.7, 4.65, WHITE, BORDER)
    card(s, 7.13, 1.35, 5.7, 0.55, RED, None)
    tb(s, 7.13, 1.35, 5.7, 0.55, [P([R('WHAT WHATSAPP GOT', 14, True, WHITE)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    ry = 2.18
    for i, (metric, label) in enumerate(L):
        yy = ry + i * 0.72
        icon(s, 0.75, yy, 0.44, '✓', ACCENT, WHITE, 15)
        tb(s, 1.35, yy - 0.06, 4.75, 0.6,
           [P([R(metric + '  ', 14, True, INK), R('·  ' + label, 11.5, False, GREY)])],
           anchor=MSO_ANCHOR.MIDDLE)
    for i, (metric, label) in enumerate(Rr):
        yy = ry + i * 0.72
        icon(s, 7.38, yy, 0.44, '✗', RED, WHITE, 15)
        tb(s, 7.98, yy - 0.06, 4.75, 0.6,
           [P([R(metric + '  ', 14, True, INK), R('·  ' + label, 11.5, False, GREY)])],
           anchor=MSO_ANCHOR.MIDDLE)
    # VS circle
    icon(s, 6.31, 3.35, 0.72, 'VS', ACCENT, WHITE, 18, line=WHITE)
    # bottom banner
    card(s, 0.5, 6.25, 12.33, 0.62, PRIMARY, None)
    tb(s, 0.5, 6.25, 12.33, 0.62, [P([R(
        'The most-distributed payment app in history processed the fewest payments.  Why?',
        17, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 3 : WHAT I DID ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'How I Approached It', 'METHOD'); add_footer(s, 3)
    stages = [
        ('\U0001F50D', 'PHASE 1', 'TEARDOWN',
         'Market + competitive + regulatory intel · 10 lenses · fact-checked', PYR[0]),
        ('\U0001F9EC', 'PHASE 2', 'DIAGNOSIS',
         '7 frameworks → one ranked causal model · adversarially red-teamed', PYR[2]),
        ('♟', 'PHASE 3', 'STRATEGY',
         '51 fixes proposed & stress-tested · RICE-ranked · scenario-modeled', PYR[0]),
    ]
    cw, aw = 3.7, 0.55
    x0 = (13.333 - (3 * cw + 2 * aw)) / 2
    fills = [C('1AA890'), SECONDARY, PRIMARY]
    for i, (gl, ph, lab, desc, _) in enumerate(stages):
        cx = x0 + i * (cw + aw)
        card(s, cx, 1.75, cw, 2.95, fills[i], None, radius=0.05)
        icon(s, cx + (cw - 0.78) / 2, 1.4, 0.78, gl, WHITE, PRIMARY, 26, line=fills[i])
        tb(s, cx + 0.2, 2.35, cw - 0.4, 2.2, [
            P([R(ph, 13, True, C('CFF3E8'))], align=PP_ALIGN.CENTER, sa=2),
            P([R(lab, 23, True, WHITE)], align=PP_ALIGN.CENTER, sa=8),
            P([R(desc, 12, False, WHITE)], align=PP_ALIGN.CENTER, ln=1.1),
        ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            ax = cx + cw + 0.06
            icon(s, ax, 2.95, 0.43, '→', ACCENT, WHITE, 18)
    # contribution strip
    card(s, x0, 5.25, 3 * cw + 2 * aw, 1.15, GREEN_TINT, BORDER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(5.25), Inches(0.14), Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    tb(s, x0 + 0.35, 5.25, 3 * cw + 2 * aw - 0.6, 1.15, [
        P([R('My contribution: ', 15, True, PRIMARY),
           R('I diagnosed the root cause, then ', 14, False, INK),
           R('proposed and stress-tested every fix', 14, True, INK),
           R(' — each solution carries a "why it works" ', 14, False, INK),
           R('and', 14, True, INK),
           R(' a "why it fails" case.', 14, False, INK)], ln=1.15),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 4 : THE MARKET ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'The Market That Should Have Been Theirs', 'MARKET'); add_footer(s, 4)
    # timeline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(2.55), Inches(11.13), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = PRIMARY; ln.line.fill.background(); ln.shadow.inherit = False
    nodes = [
        (1.3, '2016', 'UPI launches', True),
        (4.0, '2018', 'WhatsApp Pay beta', False),
        (6.66, '2020', 'Full rollout + Zero-MDR begins', True),
        (9.3, 'Aug 2022', 'P2M overtakes P2P', False),
        (11.95, '2024', '~$2.7T annual UPI volume', True),
    ]
    for nx, yr, lab, above in nodes:
        icon(s, nx - 0.14, 2.41, 0.33, '', ACCENT, WHITE, 8, line=WHITE)
        if above:
            tb(s, nx - 1.25, 1.5, 2.5, 0.95,
               [P([R(yr, 13, True, PRIMARY)], align=PP_ALIGN.CENTER, sa=1),
                P([R(lab, 11, False, INK)], align=PP_ALIGN.CENTER, ln=1.0)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        else:
            tb(s, nx - 1.25, 2.78, 2.5, 0.95,
               [P([R(yr, 13, True, PRIMARY)], align=PP_ALIGN.CENTER, sa=1),
                P([R(lab, 11, False, INK)], align=PP_ALIGN.CENTER, ln=1.0)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # stat cards
    sc = [('16B+', 'UPI transactions / month', PRIMARY, WHITE, BORDER, None),
          ('₹0', 'MDR — what merchants & apps earn per UPI txn', AMBER, AMBER_TINT, AMBER,
           'Zero-MDR (Jan 2020): the policy that quietly broke the business model.'),
          ('>50%', 'of volume now P2M (merchant) payments', PRIMARY, WHITE, BORDER, None)]
    cw = 3.9; gap = 0.3; x0 = (13.333 - (3 * cw + 2 * gap)) / 2
    for i, (num, cap, numcol, fill, bord, foot) in enumerate(sc):
        cx = x0 + i * (cw + gap)
        card(s, cx, 4.35, cw, 2.25, fill, bord, line_w=(2.0 if foot else 1.0))
        tb(s, cx + 0.15, 4.5, cw - 0.3, 0.95, [P([R(num, 46, True, numcol)], align=PP_ALIGN.CENTER)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, cx + 0.2, 5.5, cw - 0.4, 0.55, [P([R(cap, 12.5, True, INK)], align=PP_ALIGN.CENTER, ln=1.0)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        if foot:
            tb(s, cx + 0.2, 6.0, cw - 0.4, 0.55,
               [P([R(foot, 10, False, C('9C5A1E'), True)], align=PP_ALIGN.CENTER, ln=1.0)],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # ---------------- SLIDE 5 : COMPETITIVE LANDSCAPE ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'Who Won — And Why', 'MARKET'); add_footer(s, 5)
    tb(s, 0.6, 1.25, 5.8, 0.4, [P([R('UPI Market Share (%)', 13, True, PRIMARY)], align=PP_ALIGN.CENTER)],
       align=PP_ALIGN.CENTER)
    native_pie(s, 0.6, 1.6, 5.9, 5.1,
               ['PhonePe', 'Google Pay', 'Paytm', 'Others', 'WhatsApp Pay'],
               [47, 36, 7, 9, 1],
               [PRIMARY, TEAL, AMBER, GREY, ACCENT])
    fly = [('PhonePe', 'merchant lending + soundbox hardware moat', PRIMARY, BORDER),
           ('Google Pay', 'Search + Android ecosystem pull', TEAL, BORDER),
           ('Paytm', 'offline QR density + merchant network', AMBER, BORDER),
           ('WhatsApp Pay', 'distribution without a flywheel', ACCENT, RED)]
    fy = 1.55
    for i, (nm, desc, ccol, bord) in enumerate(fly):
        yy = fy + i * 1.18
        card(s, 6.9, yy, 5.93, 1.0, WHITE, bord, line_w=(2.0 if bord == RED else 1.0))
        sq = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), Inches(yy + 0.34), Inches(0.32), Inches(0.32))
        sq.fill.solid(); sq.fill.fore_color.rgb = ccol; sq.line.fill.background(); sq.shadow.inherit = False
        tb(s, 7.6, yy + 0.1, 5.1, 0.85, [
            P([R(nm, 14, True, INK)], sa=2),
            P([R(desc, 11.5, False, GREY)], ln=1.0)],
           anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 6 : ROOT CAUSE ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'Root Cause: A Business-Model Failure, Not a Product One', 'DIAGNOSIS'); add_footer(s, 6)

    def cbox(x, y, w, h, txt, fill, bord, bw=1.0, tcol=INK, fs=12, bold=True):
        card(s, x, y, w, h, fill, bord, line_w=bw)
        tb(s, x + 0.08, y, w - 0.16, h, [P([R(txt, fs, bold, tcol)], align=PP_ALIGN.CENTER, ln=1.0)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cy, ch, cw = 1.85, 1.05, 2.5
    xs = [0.5, 3.35, 6.2, 9.05]
    cbox(xs[0], cy, cw, ch, 'RBI / NPCI Zero-MDR\n(Jan 2020)', AMBER_TINT, AMBER, 2.0, C('9C5A1E'))
    cbox(xs[1], cy, cw, ch, 'UPI earns ₹0\nper transaction', WHITE, BORDER)
    cbox(xs[2], cy, cw, ch, "WhatsApp Pay's LTV\nto Meta ≈ ₹0", WHITE, PRIMARY, 3.0, PRIMARY)
    cbox(xs[3], cy, cw, ch, 'Rational\nnon-investment by Meta', WHITE, BORDER)
    for i in range(3):
        ax = xs[i] + cw + 0.02
        icon(s, ax, cy + ch / 2 - 0.18, 0.34, '→', PRIMARY, WHITE, 15)
    # fork
    f1x, f2x, fy, fw, fh = 5.6, 8.8, 4.0, 2.6, 0.85
    cbox(f1x, fy, fw, fh, 'No merchant density', WHITE, BORDER, 1.0, INK, 12.5)
    cbox(f2x, fy, fw, fh, 'No habit-formation loop', WHITE, BORDER, 1.0, INK, 12.5)
    b4cx = xs[3] + cw / 2
    arrow_line(s, b4cx, cy + ch, f1x + fw / 2, fy, PRIMARY, 2.2)
    arrow_line(s, b4cx, cy + ch, f2x + fw / 2, fy, PRIMARY, 2.2)
    # terminal
    tx, ty, tw, th = 7.2, 5.25, 2.9, 0.85
    cbox(tx, ty, tw, th, '<1% market share', RED_TINT, RED, 2.5, C('B11E2A'), 14)
    arrow_line(s, f1x + fw / 2, fy + fh, tx + tw / 2 - 0.4, ty, PRIMARY, 2.2)
    arrow_line(s, f2x + fw / 2, fy + fh, tx + tw / 2 + 0.4, ty, PRIMARY, 2.2)
    # caption
    card(s, 0.5, 6.32, 12.33, 0.5, GREEN_TINT, BORDER)
    tb(s, 0.5, 6.32, 12.33, 0.5, [P([R(
        'One policy removed the revenue.  No revenue removed the reason to invest.  Everything downstream followed.',
        13, True, PRIMARY, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 7 : FAILURE DRIVERS ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'Seven Drivers, Ranked by Weight', 'DIAGNOSIS'); add_footer(s, 7)
    drivers = ["Meta's strategic non-commitment", 'No LTV / monetization model',
               'Late regulatory clearance (~3-yr window lost)', 'No merchant-acquisition engine',
               'Habit loop never closed', 'Trust / scam perception', 'UX discoverability']
    dvals = [9.2, 8.8, 7.5, 7.1, 6.4, 5.2, 4.8]
    dcol = [PRIMARY, PRIMARY, AMBER, SECONDARY, SECONDARY, GREY, GREY]
    native_bar(s, 0.4, 1.4, 7.7, 5.25, drivers, dvals, dcol, maxval=11, number_format='0.0',
               cat_sz=10, lbl_sz=12)
    # natural experiment card
    card(s, 8.35, 1.5, 4.5, 5.05, DARK, None, radius=0.04)
    tb(s, 8.6, 1.7, 4.0, 0.4, [P([R('THE NATURAL EXPERIMENT', 14, True, ACCENT)])])
    tb(s, 8.6, 2.15, 4.0, 0.85, [P([R(
        "NPCI's 30% market-share cap was removed on Dec 31, 2024.", 12.5, False, WHITE)], ln=1.1)])
    tb(s, 8.6, 3.0, 4.0, 1.0, [P([R('+17%', 50, True, ACCENT)])])
    tb(s, 8.6, 4.05, 4.0, 0.55, [P([R("WhatsApp Pay's entire Q1-2025 growth.", 12.5, False, WHITE)], ln=1.05)])
    rl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(4.7), Inches(4.0), Inches(0.025))
    rl.fill.solid(); rl.fill.fore_color.rgb = ACCENT; rl.line.fill.background(); rl.shadow.inherit = False
    tb(s, 8.6, 4.85, 4.0, 1.5, [P([
        R('Verdict: ', 13, True, ACCENT),
        R('the cap was never the binding constraint — Meta’s will to invest was.',
          12.5, False, WHITE)], ln=1.15)])

    # ---------------- SLIDE 8 : WHERE TO PLAY / HOW TO WIN ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'The Strategic Choice', 'STRATEGY'); add_footer(s, 8)
    mx, my, mw, mh = 2.15, 1.55, 10.45, 4.7
    qw, qh = mw / 2, mh / 2
    # quadrants
    def quad(qx, qy, fill, bord, title, sub=None, tcol=INK, bw=1.2):
        card(s, qx + 0.04, qy + 0.04, qw - 0.08, qh - 0.08, fill, bord, line_w=bw)
        runs = [P([R(title, 15, True, tcol)], align=PP_ALIGN.CENTER, sa=2)]
        if sub: runs.append(P([R(sub, 11.5, False, tcol)], align=PP_ALIGN.CENTER))
        tb(s, qx + 0.2, qy + 0.2, qw - 0.4, qh - 0.4, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    quad(mx, my, C('EEF0F2'), BORDER, 'Commodity rail', 'present everywhere, earns nothing', GREY)
    quad(mx + qw, my, GREEN_TINT, ACCENT, 'COMMERCE OS', 'the target state', PRIMARY, bw=2.5)
    quad(mx, my + qh, RED_TINT, RED, 'WhatsApp Pay Today', '<1% share', C('B11E2A'), bw=2.5)
    quad(mx + qw, my + qh, C('EEF0F2'), BORDER, 'Revenue, no distribution', 'a wall with no door', GREY)
    # red dot in bottom-left quadrant
    icon(s, mx + 0.35, my + qh + qh - 0.6, 0.22, '', RED, WHITE, 8)
    # path arrow
    arrow_line(s, mx + 0.9, my + qh + qh - 0.5, mx + qw + qw * 0.55, my + qh * 0.55, ACCENT, 4.0)
    chip(s, mx + qw * 0.62, my + qh - 0.22, 4.0, 0.46,
         'The path: CTWA closed-loop checkout + merchant lending', PRIMARY, WHITE, 11, True, radius=0.5)
    # axis labels
    card(s, 0.45, my, 1.55, mh, C('EAECEF'), BORDER)
    tb(s, 0.5, my, 1.45, mh, [
        P([R('MONETIZATION', 11, True, PRIMARY)], align=PP_ALIGN.CENTER, sa=3),
        P([R('↑ MDR-bearing / LTV-positive', 10.5, True, INK)], align=PP_ALIGN.CENTER, sa=10, ln=1.0),
        P([R('Zero-MDR ↓', 10.5, False, GREY)], align=PP_ALIGN.CENTER, ln=1.0),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    card(s, mx, my + mh + 0.12, mw, 0.42, C('EAECEF'), BORDER)
    tb(s, mx, my + mh + 0.12, mw, 0.42,
       [P([R('Meta’s strategic commitment    Low  →  High', 12, True, INK)], align=PP_ALIGN.CENTER)],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 9 : 5 BIG BETS ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'The Fix: Five Sequenced Bets', 'STRATEGY'); add_footer(s, 9)
    bets = [
        ('1', 'CTWA Closed-Loop Checkout',
         'Ads → chat → pay; turns a ₹0 transaction into an ad-attributed sale.  FOUNDATION.'),
        ('2', 'Merchant-Lending Engine',
         '~7.5% take on loans to the 15M WhatsApp Business base; a merchant army with LTV attached.'),
        ('3', 'Chat-Native Red Packet',
         'WeChat hongbao analogue for viral P2P ignition.'),
        ('4', 'Anchor Merchants',
         '3–5 deals; one grocery anchor ≈ 40–160M txns / mo.'),
        ('5', '2027 Regulatory Window',
         'Cheap, high-optionality, decays fast.'),
    ]
    widths = [11.0, 10.2, 9.4, 8.6, 7.8]
    by, bh, bspace = 1.5, 0.9, 1.0
    cxC = 6.05
    band_mid = []
    for i, (num, title, desc) in enumerate(bets):
        w = widths[i]; x = cxC - w / 2; y = by + i * bspace
        band_mid.append((x, y + bh / 2, y))
        card(s, x, y, w, bh, PYR[i], None, radius=0.06)
        icon(s, x + 0.18, y + (bh - 0.6) / 2, 0.6, num, WHITE, PYR[i], 20)
        tb(s, x + 0.95, y, w - 1.2, bh, [
            P([R(title + '  ', 15, True, WHITE), R('—  ' + desc, 11.5, False, C('EAFBF3'))], ln=1.05)],
           anchor=MSO_ANCHOR.MIDDLE)
    # dependency: bet3 needs bet2 (single arrow up the left margin)
    x2, mid2, top2 = band_mid[1]
    chip(s, 0.4, 5.5, 1.62, 1.05,
         '⚠ Bet 3 needs Bet 2 first — else it spikes, then dies.',
         RED, WHITE, 9.5, True, radius=0.1)
    arrow_line(s, 1.2, 5.46, 1.0, mid2, RED, 2.8)
    # phase labels on right
    pl = [('FOUNDATION', 1.7), ('GROWTH', 3.6), ('OPTIONALITY', 5.4)]
    for lab, yy in pl:
        chip(s, 11.95, yy, 1.32, 0.42, lab, PRIMARY, WHITE, 9.5, True, radius=0.3)

    # ---------------- SLIDE 10 : REVENUE ARCHITECTURE ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, 'How ₹0 Becomes a Business', 'STRATEGY'); add_footer(s, 10)
    chip(s, 0.5, 1.25, 3.6, 0.4, 'LTV per active user / year (₹)', SECONDARY, WHITE, 11, True, radius=0.3)
    native_col(s, 1.8, 1.45, 9.7, 3.65,
               ['TODAY', '12–18 MONTHS', '24–36 MONTHS'],
               [40, 200, 1250],
               [GREY, SECONDARY, PRIMARY],
               ['₹0', '₹50–200', '₹500–2,000'],
               maxval=1550, cat_sz=12)
    caps = [('UPI P2P / P2M', 'no revenue to Meta'),
            ('CTWA closed-loop checkout', 'ad-attributed sales'),
            ('Merchant lending + insurance + credit', 'compounding LTV')]
    cw = 3.0; centers = [3.6, 6.65, 9.7]
    for (t1, t2), cxx in zip(caps, centers):
        card(s, cxx - cw / 2, 5.2, cw, 0.85, WHITE, BORDER)
        tb(s, cxx - cw / 2 + 0.1, 5.2, cw - 0.2, 0.85, [
            P([R(t1, 11.5, True, INK)], align=PP_ALIGN.CENTER, sa=1, ln=1.0),
            P([R(t2, 10.5, False, GREY)], align=PP_ALIGN.CENTER, ln=1.0)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    card(s, 0.5, 6.25, 12.33, 0.55, PRIMARY, None)
    tb(s, 0.5, 6.25, 12.33, 0.55, [P([R(
        'The LTV curve that makes Meta’s investment rational — and that no competitor can build inside an ad platform.',
        13.5, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- SLIDE 11 : PRIORITIZATION & STRESS TEST ----------------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    add_header(s, '51 Fixes, Prioritized and Stress-Tested', 'STRATEGY'); add_footer(s, 11)
    rice_cats = ['Merchant-Lending Engine', 'CTWA Closed-Loop Settlement', 'WhatsApp Business Merchant Army',
                 'Red-Packet Viral Loop', 'Grocery Anchor Partnership', 'RuPay Credit-on-UPI',
                 'In-Chat B2B Invoice + Pay', 'Soundbox-Equivalent Audio Confirm',
                 'Cross-Border / NRI Remittance', 'Insurance Micro-Products']
    rice_vals = [100, 98, 95, 88, 85, 80, 76, 72, 68, 64]
    rice_col = [PRIMARY, PRIMARY, SECONDARY, ACCENT, SECONDARY, PRIMARY, SECONDARY, SECONDARY, GREY, GREY]
    native_bar(s, 0.4, 1.35, 8.0, 5.0, rice_cats, rice_vals, rice_col, maxval=116,
               number_format='0', cat_sz=10, lbl_sz=11)
    # family legend under chart
    fam = [('Monetization', PRIMARY), ('Merchant', SECONDARY), ('Viral', ACCENT), ('Market/Other', GREY)]
    lx = 0.6
    for nm, cc in fam:
        sq = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(6.5), Inches(0.2), Inches(0.2))
        sq.fill.solid(); sq.fill.fore_color.rgb = cc; sq.line.fill.background(); sq.shadow.inherit = False
        tb(s, lx + 0.26, 6.42, 1.85, 0.35, [P([R(nm, 10.5, True, INK)])], anchor=MSO_ANCHOR.MIDDLE)
        lx += 1.95
    # right panel
    card(s, 8.6, 1.45, 4.3, 0.55, PRIMARY, None)
    tb(s, 8.6, 1.45, 4.3, 0.55, [P([R('MY STRESS-TEST METHOD', 13, True, WHITE)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    card(s, 8.6, 2.15, 4.3, 1.5, GREEN_TINT, ACCENT, line_w=1.5)
    icon(s, 8.75, 2.3, 0.4, '✓', ACCENT, WHITE, 14)
    tb(s, 9.25, 2.25, 3.55, 1.35, [
        P([R('WHY IT WORKS', 13, True, PRIMARY)], sa=3),
        P([R('Proof metric  ·  projected impact', 12, False, INK)], ln=1.1),
        P([R('the evidence the fix moves the North-Star number.', 10.5, False, GREY)], ln=1.05, sb=3)],
       anchor=MSO_ANCHOR.MIDDLE)
    card(s, 8.6, 3.8, 4.3, 1.5, RED_TINT, RED, line_w=1.5)
    icon(s, 8.75, 3.95, 0.4, '✗', RED, WHITE, 14)
    tb(s, 9.25, 3.9, 3.55, 1.35, [
        P([R('WHY IT MIGHT NOT', 13, True, C('B11E2A'))], sa=3),
        P([R('Kill-risk  ·  new risk created', 12, False, INK)], ln=1.1),
        P([R('the honest case for how each fix fails.', 10.5, False, GREY)], ln=1.05, sb=3)],
       anchor=MSO_ANCHOR.MIDDLE)
    chip(s, 8.6, 5.45, 4.3, 0.85,
         'Every solution scored on  Reach × Impact × Confidence ÷ Effort',
         SECONDARY, WHITE, 11.5, True, radius=0.12)

    # ---------------- SLIDE 12 : THREE FUTURES + THESIS ----------------
    s = prs.slides.add_slide(blank); bg(s, DARK)
    add_footer(s, 12, dark=True)
    chip(s, 5.42, 0.35, 2.5, 0.4, 'THREE FUTURES', SECONDARY, WHITE, 12, True, radius=0.4)
    scen = [
        ('BULL', '22%', ACCENT, 0.567, 2.95,
         ['Meta commits $500M+', '12–15% UPI share', '₹2,500–4,000cr revenue'],
         'India treated as a commerce OS.'),
        ('BASE', '50%', AMBER, 4.767, 3.15,
         ['CTWA works, lending launches', 'daily P2P never fully converts', '4–6% UPI share'],
         'Half-commitment — money flows, not enough.'),
        ('BEAR', '28%', RED, 8.967, 2.95,
         ['Meta stays passive', '<1.5% share', 'a second era of under-investment'],
         'A replay of Dec 2024.'),
    ]
    for name, prob, col, x, h, metrics, cond in scen:
        y = 0.95 if h > 3.0 else 1.05
        card(s, x, y, 3.8, h, CARD_DK, col, line_w=2.5)
        icon(s, x + (3.8 - 0.82) / 2, y - 0.2, 0.82, prob, col, WHITE, 19, line=DARK)
        runs = [P([R(name, 17, True, col)], align=PP_ALIGN.CENTER, sa=8)]
        for m in metrics:
            runs.append(P([R(m, 12.5, False, WHITE)], align=PP_ALIGN.CENTER, sa=4, ln=1.0))
        tb(s, x + 0.2, y + 0.68, 3.4, h - 1.35, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        rl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.4), Inches(y + h - 0.95), Inches(3.0), Inches(0.02))
        rl.fill.solid(); rl.fill.fore_color.rgb = col; rl.line.fill.background(); rl.shadow.inherit = False
        tb(s, x + 0.2, y + h - 0.85, 3.4, 0.75,
           [P([R('Condition: ', 10.5, True, col), R(cond, 10.5, False, C('C9D1D6'))],
              align=PP_ALIGN.CENTER, ln=1.0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # thesis
    tb(s, 1.0, 4.12, 11.33, 2.5, [
        P([R('The real question for Meta isn’t', 20, False, WHITE)], align=PP_ALIGN.CENTER, sa=3),
        P([R('“How do we grow UPI share?”', 28, True, ACCENT)], align=PP_ALIGN.CENTER, sa=5),
        P([R('It’s', 18, False, WHITE)], align=PP_ALIGN.CENTER, sa=3),
        P([R('“Do we build a commerce OS on the world’s most-used messaging app — or not?”',
             24, True, ACCENT)], align=PP_ALIGN.CENTER, sa=7, ln=1.05),
        P([R('Every other variable is downstream of that one decision.', 14, False, GREY, True)],
          align=PP_ALIGN.CENTER),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    tb(s, 0.5, 6.95, 12.33, 0.3, [P([R(
        'Sources: NPCI  ·  RBI payment-system reports  ·  PhonePe FY25 filings  ·  Meta earnings calls  ·  industry data',
        9, False, GREY)], align=PP_ALIGN.CENTER)], align=PP_ALIGN.CENTER)

    # ---------------- save ----------
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'WhatsApp_Pay_Case_Study.pptx')
    prs.save(out)
    print('Saved:', out)
    print('Slides:', len(prs.slides._sldIdLst))
    titles = ['Title / Hero', 'The Paradox', 'How I Approached It (Method)', 'The Market',
              'Competitive Landscape (pie)', 'Root Cause Tree', 'Failure Drivers (bar) + Natural Experiment',
              'Where to Play / How to Win (2x2)', 'Five Big Bets (pyramid)', 'Revenue Architecture (columns)',
              'Prioritization & Stress Test (bar)', 'Three Futures + Thesis']
    for i, t in enumerate(titles, 1):
        print('  %2d. %s' % (i, t))


if __name__ == '__main__':
    build()
