"""
WhatsApp Pay India — Product Strategy Case Study
Master deck builder. Dark neon-green theme. python-pptx.
Rebuilds the whole .pptx from scratch each run (idempotent).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── PALETTE ──────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0A, 0x0A)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
GREEN_MID = RGBColor(0x1B, 0x5E, 0x20)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xB0, 0xBE, 0xC5)
GREY_DK   = RGBColor(0x6C, 0x77, 0x7E)
CARD      = RGBColor(0x16, 0x16, 0x16)
PANEL     = RGBColor(0x0D, 0x0D, 0x0D)
ORANGE    = RGBColor(0xFF, 0x6F, 0x00)
RED       = RGBColor(0xD3, 0x2F, 0x2F)
WM        = RGBColor(0x0F, 0x1F, 0x0F)
# competitor colors
C_PHONEPE = RGBColor(0x00, 0x69, 0x5C)
C_GPAY    = RGBColor(0x02, 0x88, 0xD1)
C_PAYTM   = RGBColor(0xF5, 0x7C, 0x00)
C_OTHERS  = RGBColor(0x54, 0x6E, 0x7A)
C_WA      = RGBColor(0x00, 0xE6, 0x76)

A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSC = "http://schemas.openxmlformats.org/drawingml/2006/chart"

def hx(c):  # RGBColor -> 'RRGGBB'
    return "%02X%02X%02X" % (c[0], c[1], c[2])

# ── LOW-LEVEL HELPERS ────────────────────────────────────
def flatten(shape):
    sp = shape._element
    for child in list(sp):
        if child.tag == qn('p:style'):
            sp.remove(child)

def glow_run(run, hexcolor="00E676", rad=88900, alpha=55000):
    rPr = run._r.get_or_add_rPr()
    eff = parse_xml(
        f'<a:effectLst xmlns:a="{A}">'
        f'<a:glow rad="{rad}"><a:srgbClr val="{hexcolor}"><a:alpha val="{alpha}"/></a:srgbClr></a:glow>'
        f'</a:effectLst>')
    after = [qn('a:highlight'), qn('a:uLnTx'), qn('a:uLn'), qn('a:uFillTx'), qn('a:uFill'),
             qn('a:latin'), qn('a:ea'), qn('a:cs'), qn('a:sym'), qn('a:hlinkClick'),
             qn('a:hlinkMouseOver'), qn('a:rtl'), qn('a:extLst')]
    for child in rPr:
        if child.tag in after:
            child.addprevious(eff); return
    rPr.append(eff)

def glow_shape(shape, hexcolor="00E676", rad=76200, alpha=50000):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    spPr.append(parse_xml(
        f'<a:effectLst xmlns:a="{A}">'
        f'<a:glow rad="{rad}"><a:srgbClr val="{hexcolor}"><a:alpha val="{alpha}"/></a:srgbClr></a:glow>'
        f'</a:effectLst>'))

def textbox(slide, text, x, y, w, h, *, size=12, color=WHITE, bold=False, italic=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, glow=None, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    if glow: glow_run(r, *glow)
    return tb

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    flatten(s)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s

def rrect(slide, x, y, w, h, fill, line=None, lw=Pt(1.5), radius=0.1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    flatten(s)
    try: s.adjustments[0] = radius
    except Exception: pass
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s

def oval(slide, cx, cy, r, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    flatten(s)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s

def vline(slide, x, y1, y2, color, width=Pt(1), dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    ln.line.color.rgb = color; ln.line.width = width
    ln.shadow.inherit = False
    if dash:
        lnEl = ln.line._get_or_add_ln()
        lnEl.append(parse_xml(f'<a:prstDash xmlns:a="{A}" val="{dash}"/>'))
    return ln

def arrow(slide, x1, y1, x2, y2, color=GREEN, width=Pt(2.25), big=False, glow=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width; cn.shadow.inherit = False
    sz = "lg" if big else "med"
    cn.line._get_or_add_ln().append(
        parse_xml(f'<a:tailEnd xmlns:a="{A}" type="triangle" w="{sz}" len="{sz}"/>'))
    if glow:
        sp = cn._element.spPr
        for el in sp.findall(qn('a:effectLst')): sp.remove(el)
        sp.append(parse_xml(
            f'<a:effectLst xmlns:a="{A}"><a:glow rad="63500">'
            f'<a:srgbClr val="{hx(color)}"><a:alpha val="45000"/></a:srgbClr></a:glow></a:effectLst>'))
    return cn

def fill_tf(tf, text, size, color, *, bold=False, italic=False, align=PP_ALIGN.CENTER, font="Calibri"):
    """Fill a text frame, splitting on \\n into centered paragraphs."""
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, font
        r.font.color.rgb = color

def para(tf, text, size, color, *, bold=False, italic=False, before=0, align=PP_ALIGN.LEFT, first=False, font="Calibri"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, font
    r.font.color.rgb = color
    return r

def pill(slide, text, x, y, w, h, *, fill=GREEN_MID, line=GREEN, txt=GREEN, size=9, bold=True):
    s = rrect(slide, x, y, w, h, fill, line, Pt(1), radius=0.5)
    tf = s.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, "Calibri"
    r.font.color.rgb = txt
    return s

def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def label_pill(slide, text):
    pill(slide, text, Inches(11.45), Inches(0.18), Inches(1.55), Inches(0.34))

def title(slide, text, size=28):
    textbox(slide, text, Inches(0.4), Inches(0.3), Inches(10.6), Inches(0.62),
            size=size, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Inches(0.4), Inches(0.92), Inches(12.53), Emu(26000), GREEN)

def footer(slide, n):
    textbox(slide, "WhatsApp Pay India — Product Strategy Case Study", Inches(0.4),
            Inches(7.18), Inches(8), Inches(0.25), size=9, color=GREY_DK)
    textbox(slide, f"{n} / 14", Inches(11.9), Inches(7.18), Inches(1.03), Inches(0.25),
            size=9, color=GREY_DK, align=PP_ALIGN.RIGHT)

def banner(slide, text, y, color=GREEN):
    rect(slide, Inches(0.4), y, Inches(12.53), Inches(0.58), RGBColor(0x11, 0x11, 0x11),
         line=RGBColor(0x1B, 0x5E, 0x20), lw=Pt(0.75))
    textbox(slide, text, Inches(0.6), y, Inches(12.13), Inches(0.58), size=14, color=color,
            bold=True, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── CHART HELPERS ────────────────────────────────────────
def chart_transparent(chart):
    cs = chart._chartSpace
    spxml = f'<c:spPr xmlns:c="{NSC}" xmlns:a="{A}"><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'
    for el in cs.findall(qn('c:spPr')): cs.remove(el)
    chart_el = cs.find(qn('c:chart'))
    chart_el.addnext(parse_xml(spxml))
    pa = chart_el.find(qn('c:plotArea'))
    for el in pa.findall(qn('c:spPr')): pa.remove(el)
    pa.append(parse_xml(spxml))

def set_hole_size(chart, pct=58):
    dough = chart._chartSpace.find(qn('c:chart')).find(qn('c:plotArea')).find(qn('c:doughnutChart'))
    for el in dough.findall(qn('c:holeSize')): dough.remove(el)
    dough.append(parse_xml(f'<c:holeSize xmlns:c="{NSC}" val="{pct}"/>'))

def add_donut(slide, x, y, w, h, cats, vals, colors):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('Share', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd)
    ch = gf.chart
    chart_transparent(ch)
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10); ch.legend.font.color.rgb = GREY; ch.legend.font.name = "Calibri"
    plot = ch.plots[0]; ser = plot.series[0]
    for i, c in enumerate(colors):
        pt = ser.points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
        pt.format.line.color.rgb = BG; pt.format.line.width = Pt(2)
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.show_category_name = False; dl.show_percentage = False
    dl.show_legend_key = False; dl.show_series_name = False
    dl.number_format = '0"%"'; dl.number_format_is_linked = False
    dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = WHITE; dl.font.name = "Calibri"
    set_hole_size(ch, 58)
    return ch

def add_hbar(slide, x, y, w, h, cats, vals, colors, vmax=50, num_fmt='0.0"%"'):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('v', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd)
    ch = gf.chart
    chart_transparent(ch)
    ch.has_title = False; ch.has_legend = False
    plot = ch.plots[0]; plot.gap_width = 55; ser = plot.series[0]
    for i, c in enumerate(colors):
        ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb = c
        ser.points[i].format.line.fill.background()
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt; dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10); dl.font.bold = True; dl.font.color.rgb = WHITE; dl.font.name = "Calibri"
    cax = ch.category_axis
    cax.tick_labels.font.size = Pt(10.5); cax.tick_labels.font.color.rgb = WHITE
    cax.tick_labels.font.name = "Calibri"
    cax.format.line.color.rgb = GREY_DK
    cax.has_major_gridlines = False
    vax = ch.value_axis
    vax.visible = False
    vax.has_major_gridlines = False; vax.has_minor_gridlines = False
    vax.maximum_scale = vmax; vax.minimum_scale = 0
    return ch

def add_column(slide, x, y, w, h, cats, vals, colors, vmax=None, gap=80):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('v', vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    ch = gf.chart
    chart_transparent(ch)
    ch.has_title = False; ch.has_legend = False
    plot = ch.plots[0]; plot.gap_width = gap; ser = plot.series[0]
    for i, c in enumerate(colors):
        ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb = c
        ser.points[i].format.line.fill.background()
    plot.has_data_labels = False
    vax = ch.value_axis; vax.visible = False
    vax.has_major_gridlines = False; vax.has_minor_gridlines = False
    if vmax: vax.maximum_scale = vmax; vax.minimum_scale = 0
    cax = ch.category_axis; cax.visible = False; cax.has_major_gridlines = False
    return ch

def add_stacked_column(slide, x, y, w, h, cats, series, colors, vmax=None, gap=80):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series: cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, w, h, cd)
    ch = gf.chart
    chart_transparent(ch)
    ch.has_title = False; ch.has_legend = False
    plot = ch.plots[0]; plot.gap_width = gap
    for si, c in enumerate(colors):
        ser = plot.series[si]
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
        ser.format.line.color.rgb = BG; ser.format.line.width = Pt(1)
    plot.has_data_labels = False
    vax = ch.value_axis; vax.visible = False
    vax.has_major_gridlines = False; vax.has_minor_gridlines = False
    if vmax: vax.maximum_scale = vmax; vax.minimum_scale = 0
    cax = ch.category_axis; cax.visible = False; cax.has_major_gridlines = False
    return ch

# ═════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════
def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    SW = prs.slide_width
    textbox(s, "₹", Inches(0), Inches(0.2), SW, Inches(7.1), size=360, color=WM, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, "A PRODUCT STRATEGY CASE STUDY", Inches(0), Inches(1.95), SW, Inches(0.4),
            size=12, color=GREY, align=PP_ALIGN.CENTER)
    hl = s.shapes.add_textbox(Inches(0.4), Inches(2.45), Inches(12.53), Inches(1.7))
    htf = hl.text_frame; htf.word_wrap = True; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    htf.margin_left = htf.margin_right = htf.margin_top = htf.margin_bottom = 0
    p1 = htf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = "500 Million Users.  ₹0 Revenue."
    r1.font.size, r1.font.bold, r1.font.name = Pt(41), True, "Calibri"; r1.font.color.rgb = WHITE
    p2 = htf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = "Why WhatsApp Pay Lost India."
    r2.font.size, r2.font.bold, r2.font.name = Pt(41), True, "Calibri"; r2.font.color.rgb = GREEN
    glow_run(r2, "00E676", rad=88900, alpha=60000)
    rect(s, (SW - Inches(2.4)) // 2, Inches(4.28), Inches(2.4), Emu(28000), GREEN)
    textbox(s, "Competitive Teardown   →   Root-Cause Diagnosis   →   Rejuvenation Strategy",
            Inches(0), Inches(4.52), SW, Inches(0.4), size=13, color=GREY, align=PP_ALIGN.CENTER)
    labels = ["3 PHASES", "84,000+ WORDS", "51 SOLUTIONS", "RICE-PRIORITIZED"]
    pw, ph, gap = Inches(1.85), Inches(0.34), Inches(0.2)
    total = pw * len(labels) + gap * (len(labels) - 1); px = (SW - total) // 2
    for lab in labels:
        pill(s, lab, px, Inches(6.35), pw, ph, size=9.5); px += pw + gap
    textbox(s, "Fahad Ahmed  ·  PM Portfolio", Inches(0.4), Inches(7.02), Inches(5), Inches(0.32),
            size=10, color=GREY_DK)
    textbox(s, "fahadahmed982004@gmail.com", Inches(7.93), Inches(7.02), Inches(5), Inches(0.32),
            size=10, color=GREY_DK, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════
# SLIDE 2 — THE MARKET
# ═════════════════════════════════════════════════════════
def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "MARKET")
    title(s, "The Market That Should Have Been Theirs", size=28)
    cy, ch, cw, cg = Inches(1.18), Inches(2.15), Inches(4.04), Inches(0.2); cx0 = Inches(0.4)
    cards = [("16B+", "UPI transactions / month", GREEN, None),
             ("$2.7T", "Annual UPI volume (2024)", GREEN, None),
             (">50%", "of volume now P2M (merchant)", ORANGE,
              "Zero-MDR (Jan 2020): the policy that quietly broke the business model.")]
    for i, (big, sub, border, caption) in enumerate(cards):
        cx = cx0 + i * (cw + cg)
        rrect(s, cx, cy, cw, ch, CARD, border, Pt(2), radius=0.06)
        fr = s.shapes.add_textbox(cx + Inches(0.18), cy + Inches(0.12), cw - Inches(0.36), ch - Inches(0.24))
        tf = fr.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        pa = tf.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
        ra = pa.add_run(); ra.text = big
        ra.font.size, ra.font.bold, ra.font.name = Pt(46), True, "Calibri"; ra.font.color.rgb = GREEN
        glow_run(ra, "00E676", rad=63500, alpha=50000)
        pb = tf.add_paragraph(); pb.alignment = PP_ALIGN.CENTER; pb.space_before = Pt(8)
        rb = pb.add_run(); rb.text = sub
        rb.font.size, rb.font.name = Pt(12.5), "Calibri"; rb.font.color.rgb = WHITE
        if caption:
            pc = tf.add_paragraph(); pc.alignment = PP_ALIGN.CENTER; pc.space_before = Pt(8)
            rc = pc.add_run(); rc.text = caption
            rc.font.size, rc.font.italic, rc.font.name = Pt(9), True, "Calibri"; rc.font.color.rgb = ORANGE
    panel_y = Inches(3.6)
    rect(s, Inches(0.4), panel_y, Inches(12.53), Inches(3.55), PANEL)
    textbox(s, "TEN YEARS, ONE MISSED MARKET", Inches(0.65), panel_y + Inches(0.18),
            Inches(8), Inches(0.3), size=11, color=GREY_DK, bold=True)
    line_y = panel_y + Inches(2.3)
    tl_x0, tl_x1 = Inches(1.0), Inches(12.3); tl_w = tl_x1 - tl_x0
    rect(s, tl_x0, line_y - Emu(16000), tl_w, Emu(32000), GREEN_MID)
    events = [("2016", "UPI launches", GREY, None),
              ("2018", "WhatsApp Pay beta", GREEN, None),
              ("Jan 2020", "Zero-MDR begins", ORANGE, "₹0 per transaction — the\nrevenue model breaks"),
              ("Dec 2020", "30% TPAP cap imposed", ORANGE, "No app may hold >30%\nof UPI volume"),
              ("Dec 2024", "NPCI cap removed", GREEN, None),
              ("Jun 2026", "Kunal Shah → WhatsApp", GREEN, None)]
    n = len(events)
    for i, (yr, ev, col, note) in enumerate(events):
        frac = 0.04 + i * (0.92 / (n - 1)); ex = tl_x0 + int(tl_w * frac)
        d = oval(s, ex, line_y, Inches(0.11), col); glow_shape(d, hx(col), rad=63500, alpha=45000)
        textbox(s, yr, ex - Inches(0.9), line_y - Inches(0.5), Inches(1.8), Inches(0.3),
                size=10, color=col, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, ev, ex - Inches(1.0), line_y + Inches(0.18), Inches(2.0), Inches(0.5),
                size=9.5, color=(ORANGE if col == ORANGE else GREY), align=PP_ALIGN.CENTER)
        if note:
            bw, bh = Inches(1.95), Inches(0.66); bx = ex - bw // 2; by = line_y - Inches(1.42)
            rrect(s, bx, by, bw, bh, RGBColor(0x1F, 0x1F, 0x00), ORANGE, Pt(1), radius=0.12)
            cb = s.shapes.add_textbox(bx + Inches(0.08), by, bw - Inches(0.16), bh)
            ctf = cb.text_frame; ctf.word_wrap = True; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
            for j, ln in enumerate(note.split("\n")):
                pp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
                rr = pp.add_run(); rr.text = ln
                rr.font.size, rr.font.italic, rr.font.name = Pt(8.5), True, "Calibri"; rr.font.color.rgb = ORANGE
    footer(s, 2)

# ═════════════════════════════════════════════════════════
# SLIDE 3 — WHO WON AND WHY
# ═════════════════════════════════════════════════════════
def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "MARKET")
    title(s, "Who Won — And Why", size=28)

    # LEFT — donut (market share by value)
    textbox(s, "UPI MARKET SHARE — BY VALUE (%)", Inches(0.5), Inches(1.02), Inches(5.6), Inches(0.3),
            size=11, color=GREY_DK, bold=True)
    add_donut(s, Inches(0.45), Inches(1.3), Inches(5.7), Inches(4.95),
              ["PhonePe", "Google Pay", "Paytm", "Others", "WhatsApp Pay"],
              [47, 36, 7, 9, 1],
              [C_PHONEPE, C_GPAY, C_PAYTM, C_OTHERS, C_WA])

    # RIGHT — competitor cards
    rx, rw = Inches(6.45), Inches(6.48)
    comp = [("PhonePe", "Merchant lending + Soundbox hardware moat", C_PHONEPE, C_PHONEPE, False),
            ("Google Pay", "Search + Android ecosystem pull", C_GPAY, C_GPAY, False),
            ("Paytm", "Offline QR density + merchant network", C_PAYTM, C_PAYTM, False),
            ("WhatsApp Pay", "Distribution without a flywheel — 500M users, no habit", C_WA, RED, True)]
    cy = Inches(1.18); chh = Inches(0.74); cgap = Inches(0.11)
    for name, desc, dot, border, outlier in comp:
        rrect(s, rx, cy, rw, chh, CARD, border, Pt(2.25) if outlier else Pt(1.25), radius=0.14)
        oval(s, rx + Inches(0.34), cy + chh // 2, Inches(0.095), dot)
        fr = s.shapes.add_textbox(rx + Inches(0.62), cy, rw - Inches(0.8), chh)
        tf = fr.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = name
        r1.font.size, r1.font.bold, r1.font.name = Pt(13.5), True, "Calibri"; r1.font.color.rgb = WHITE
        if outlier:
            rb = p1.add_run(); rb.text = "    ✕ the outlier"
            rb.font.size, rb.font.bold, rb.font.name = Pt(9.5), True, "Calibri"; rb.font.color.rgb = RED
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = desc
        r2.font.size, r2.font.name = Pt(10.5), "Calibri"; r2.font.color.rgb = GREY
        cy += chh + cgap

    # RIGHT lower — transaction volume bar
    textbox(s, "AND BY TRANSACTION VOLUME (%) — WHO ACTUALLY GETS USED", Inches(6.5), Inches(4.5),
            Inches(6.4), Inches(0.3), size=11, color=GREY_DK, bold=True)
    add_hbar(s, Inches(6.2), Inches(4.78), Inches(6.7), Inches(1.5),
             ["WhatsApp Pay", "Paytm", "Others", "Google Pay", "PhonePe"],
             [0.6, 6, 8.4, 37, 48],
             [C_WA, C_PAYTM, C_OTHERS, C_GPAY, C_PHONEPE], vmax=55)

    banner(s, "The most-distributed payment app in history processed the fewest payments.  Why?",
           Inches(6.52))
    footer(s, 3)

# ═════════════════════════════════════════════════════════
# SLIDE 4 — THE COMPETITIVE TIMELINE (swim lanes)
# ═════════════════════════════════════════════════════════
def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "MARKET")
    title(s, "A Race WhatsApp Was Always Losing", size=28)

    Y0, Y1 = 2016, 2026
    lane_x0, lane_x1 = Inches(2.0), Inches(12.5)
    lane_w = lane_x1 - lane_x0
    def yx(year):
        return lane_x0 + int(lane_w * (year - Y0) / (Y1 - Y0))

    lanes_top = Inches(1.62)
    lane_h = Inches(1.18)
    n_lanes = 4
    lanes = [
        ("PhonePe", C_PHONEPE, False, [
            (2016, "Launch"), (2017, "50M users"), (2018, "Market leader"),
            (2020, "Zero-MDR\nabsorbed"), (2022, "Soundbox"), (2024, "47% share")]),
        ("Google Pay", C_GPAY, False, [
            (2017, "Tez launch"), (2018, "100M users"), (2019, "UPI dominant"),
            (2023, "Android\nlock-in"), (2024, "36% share")]),
        ("Paytm", C_PAYTM, False, [
            (2016, "Demonetization\nsurge"), (2018, "QR network"), (2020, "MDR loss\nhurts"),
            (2022, "IPO struggles"), (2024, "7% share ↓")]),
        ("WhatsApp Pay", C_WA, True, [
            (2018, "Beta", GREEN), (2020, "Approved,\ncapped 20M", ORANGE),
            (2022, "Cap\nextended", ORANGE), (2024, "Cap\nremoved", GREEN),
            (2025, "+17%", GREEN), (2026, "Kunal Shah", GREEN)]),
    ]

    # lane band backgrounds (alternating) + year gridlines
    for i in range(n_lanes):
        ly = lanes_top + i * lane_h
        if i % 2 == 0:
            rect(s, lane_x0 - Inches(1.3), ly, lane_w + Inches(1.3), lane_h, RGBColor(0x0E, 0x0E, 0x0E))
    bands_bottom = lanes_top + n_lanes * lane_h
    for yr in range(Y0, Y1 + 1):
        gx = yx(yr)
        vline(s, gx, lanes_top, bands_bottom, RGBColor(0x1A, 0x1A, 0x1A), width=Pt(0.75))
        textbox(s, str(yr), gx - Inches(0.35), bands_bottom + Inches(0.04), Inches(0.7), Inches(0.25),
                size=8.5, color=GREY_DK, align=PP_ALIGN.CENTER)
    # divider separating lane labels from the timeline
    vline(s, lane_x0 - Inches(0.4), lanes_top, bands_bottom, RGBColor(0x22, 0x22, 0x22), width=Pt(0.75))

    # annotation vertical lines (regulatory wall)
    for yr in (2020.0, 2020.92):
        vline(s, yx(yr), lanes_top, bands_bottom, ORANGE, width=Pt(1.25), dash="dash")

    # lane baselines, labels, bubbles
    for i, (name, color, wa, events) in enumerate(lanes):
        ly = lanes_top + i * lane_h
        base_y = ly + lane_h // 2
        # baseline
        bl = rect(s, lane_x0, base_y - Emu(11000), lane_w, Emu(22000), color)
        if wa:
            bl.line.fill.background()  # keep solid; WA recovery shown via colored bubbles
        # left label
        textbox(s, name, Inches(0.22), base_y - Inches(0.2), Inches(1.25), Inches(0.4),
                size=12, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # bubbles (alternate above / below)
        for k, ev in enumerate(events):
            yr, lab = ev[0], ev[1]
            bcol = ev[2] if len(ev) > 2 else color
            cx = yx(yr)
            up = (k % 2 == 0)
            cyc = base_y - Inches(0.30) if up else base_y + Inches(0.30)
            # tick from baseline to bubble
            vline(s, cx, base_y, cyc, bcol, width=Pt(1))
            txt = lab
            w = Inches(max(0.78, min(1.5, 0.34 + 0.082 * max(len(s2) for s2 in txt.split("\n")))))
            h = Inches(0.46) if "\n" in txt else Inches(0.32)
            # clamp bubble within the timeline (keep clear of left labels & right edge)
            bx = cx - w // 2
            bx = max(Inches(1.52), min(bx, Inches(12.91) - w))
            bub = rrect(s, bx, cyc - h // 2, w, h, CARD, bcol, Pt(1.25), radius=0.22)
            tf = bub.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = 0
            for j, lpart in enumerate(txt.split("\n")):
                pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
                rr = pp.add_run(); rr.text = lpart
                rr.font.size, rr.font.bold, rr.font.name = Pt(8), True, "Calibri"
                rr.font.color.rgb = WHITE

    # top annotation callout
    bx, by, bw, bh = yx(2020.4) - Inches(1.7), Inches(0.98), Inches(3.4), Inches(0.5)
    rrect(s, bx, by, bw, bh, RGBColor(0x1F, 0x12, 0x00), ORANGE, Pt(1), radius=0.1)
    cb = s.shapes.add_textbox(bx + Inches(0.1), by, bw - Inches(0.2), bh)
    ctf = cb.text_frame; ctf.word_wrap = True; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
    for j, ln in enumerate(["JAN 2020 · Zero-MDR — ₹0 per transaction, model breaks",
                            "DEC 2020 · 30% TPAP cap (extended repeatedly to 2026)"]):
        pp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = ln
        rr.font.size, rr.font.bold, rr.font.name = Pt(8.5), True, "Calibri"; rr.font.color.rgb = ORANGE

    footer(s, 4)

# ── SHARED: flow box & quadrant ──────────────────────────
def flow_box(s, x, y, w, h, text, fill, border, txt=WHITE, lw=Pt(1.75), size=12, glow=False, bold=True):
    b = rrect(s, x, y, w, h, fill, border, lw, radius=0.12)
    if glow: glow_shape(b, hx(border), rad=63500, alpha=45000)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(6); tf.margin_top = tf.margin_bottom = Pt(2)
    fill_tf(tf, text, size, txt, bold=bold)
    return b

def quad(s, x, y, w, h, fill, border, t_title, t_color, t_size, sub, sub_color, glow=False):
    b = rrect(s, x, y, w, h, fill, border, Pt(2.25), radius=0.04)
    if glow: glow_shape(b, hx(border), rad=88900, alpha=40000)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(12); tf.margin_top = Pt(14)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t_title
    r.font.size, r.font.bold, r.font.name = Pt(t_size), True, "Calibri"; r.font.color.rgb = t_color
    if glow: glow_run(r, hx(t_color), rad=76200, alpha=55000)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(5)
    r2 = p2.add_run(); r2.text = sub
    r2.font.size, r2.font.italic, r2.font.name = Pt(11), True, "Calibri"; r2.font.color.rgb = sub_color
    return b

# ═════════════════════════════════════════════════════════
# SLIDE 5 — ROOT CAUSE
# ═════════════════════════════════════════════════════════
def slide5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "DIAGNOSIS")
    title(s, "Root Cause: A Business-Model Failure, Not a Product One", size=24)

    GRYBORD = RGBColor(0x33, 0x33, 0x33)
    bw, bh = Inches(2.45), Inches(0.62)
    xs = [Inches(0.45), Inches(3.76), Inches(7.07), Inches(10.38)]
    r1y = Inches(1.12)
    b3y, b3h = Inches(1.05), Inches(0.76)
    flow_box(s, xs[0], r1y, bw, bh, "RBI / NPCI Zero-MDR\n(Jan 2020)", RGBColor(0x1F, 0x10, 0x00), ORANGE, WHITE, size=11.5)
    flow_box(s, xs[1], r1y, bw, bh, "UPI earns ₹0\nper transaction", RGBColor(0x11, 0x11, 0x11), GRYBORD, WHITE, size=11.5)
    flow_box(s, xs[2], b3y, bw, b3h, "WhatsApp Pay's LTV\nto Meta ≈ ₹0", RGBColor(0x00, 0x22, 0x00), GREEN, WHITE, size=12, glow=True)
    flow_box(s, xs[3], r1y, bw, bh, "Rational non-\ninvestment by Meta", RGBColor(0x11, 0x11, 0x11), GRYBORD, WHITE, size=11.5)
    midy = r1y + bh // 2
    for a, b in [(0, 1), (1, 2), (2, 3)]:
        arrow(s, xs[a] + bw, midy, xs[b], midy, GREEN, Pt(2.25))

    r2y, ch = Inches(2.16), Inches(0.55)
    flow_box(s, xs[2], r2y, bw, ch, "No merchant density", CARD, GRYBORD, GREY, size=11.5)
    flow_box(s, xs[3], r2y, bw, ch, "No habit-formation loop", CARD, GRYBORD, GREY, size=11.5)
    cx3, cx4 = xs[2] + bw // 2, xs[3] + bw // 2
    arrow(s, cx3, b3y + b3h, cx3, r2y, GREEN, Pt(2.25))
    arrow(s, cx4, r1y + bh, cx4, r2y, GREEN, Pt(2.25))
    fmid = (cx3 + cx4) // 2
    fw, fy, fh = Inches(3.2), Inches(3.0), Inches(0.6)
    flow_box(s, fmid - fw // 2, fy, fw, fh, "< 1% market share", RGBColor(0x1A, 0x00, 0x00), RED, RED, size=16)
    arrow(s, cx3, r2y + ch, fmid - Inches(0.5), fy, GREEN, Pt(2.25))
    arrow(s, cx4, r2y + ch, fmid + Inches(0.5), fy, GREEN, Pt(2.25))

    # left insight
    rect(s, Inches(0.55), Inches(2.28), Emu(42000), Inches(1.05), GREEN)
    ins = s.shapes.add_textbox(Inches(0.78), Inches(2.28), Inches(5.3), Inches(1.05))
    itf = ins.text_frame; itf.word_wrap = True; itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    itf.margin_left = itf.margin_top = itf.margin_bottom = 0
    fill_tf(itf, "Each link is individually rational.\nChained together, <1% share becomes the\nequilibrium — not an accident.",
            12, GREY, italic=True, align=PP_ALIGN.LEFT)

    # ── bottom: ranked bar + natural-experiment card ──
    textbox(s, "SEVEN DRIVERS OF FAILURE — RANKED BY WEIGHT", Inches(0.45), Inches(3.92),
            Inches(8.5), Inches(0.3), size=12, color=GREY_DK, bold=True)
    add_hbar(s, Inches(0.25), Inches(4.2), Inches(8.7), Inches(2.12),
             ["UX discoverability", "Trust / scam perception", "Habit loop never closed",
              "No merchant-acquisition engine", "Late regulatory clearance",
              "No LTV / monetization model", "Meta's strategic non-commitment"],
             [4.8, 5.2, 6.4, 7.1, 7.5, 8.8, 9.2],
             [C_OTHERS, C_OTHERS, GREEN, GREEN, ORANGE, GREEN, GREEN], vmax=10, num_fmt='0.0')

    # natural experiment card
    nx, ny, nw, nh = Inches(9.15), Inches(4.02), Inches(3.78), Inches(2.5)
    rrect(s, nx, ny, nw, nh, RGBColor(0x14, 0x14, 0x14), GREEN, Pt(1.5), radius=0.06)
    nb = s.shapes.add_textbox(nx + Inches(0.18), ny + Inches(0.12), nw - Inches(0.36), nh - Inches(0.24))
    ntf = nb.text_frame; ntf.word_wrap = True; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = 0
    def addp(tf, text, size, color, bold=False, italic=False, before=0, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if before: p.space_before = Pt(before)
        r = p.add_run(); r.text = text
        r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, "Calibri"
        r.font.color.rgb = color
    addp(ntf, "THE NATURAL EXPERIMENT", 11.5, GREEN, bold=True, first=True)
    addp(ntf, "NPCI 30% cap removed — Dec 2024.", 11, WHITE, before=8)
    addp(ntf, "WhatsApp Pay then grew +17% in Q1 2025.", 11, WHITE, before=4)
    addp(ntf, "VERDICT", 9.5, GREY_DK, bold=True, before=10)
    addp(ntf, "The cap was never the binding constraint — Meta's will to invest was.",
         11.5, GREEN, italic=True, before=3)

    banner(s, "One policy removed the revenue.   No revenue removed the reason to invest.   Everything downstream followed.",
           Inches(6.5))
    footer(s, 5)

# ═════════════════════════════════════════════════════════
# SLIDE 6 — THE STRATEGIC CHOICE (2×2 matrix)
# ═════════════════════════════════════════════════════════
def slide6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "STRATEGY")
    title(s, "The Strategic Choice", size=28)

    xL, xR = Inches(1.95), Inches(12.9)
    yT, yB = Inches(1.5), Inches(6.18)
    gap = Inches(0.13)
    qw = (xR - xL - gap) // 2
    qh = (yB - yT - gap) // 2
    xRcol = xL + qw + gap
    yBrow = yT + qh + gap

    GRYBORD = RGBColor(0x44, 0x44, 0x44)
    quad(s, xL, yT, qw, qh, RGBColor(0x11, 0x11, 0x11), GRYBORD,
         "Revenue, No Distribution", WHITE, 15, "a wall with no door", GREY)
    quad(s, xRcol, yT, qw, qh, RGBColor(0x0A, 0x2A, 0x0A), GREEN,
         "COMMERCE OS", GREEN, 22, "the target state", WHITE, glow=True)
    quad(s, xL, yBrow, qw, qh, RGBColor(0x11, 0x11, 0x11), GRYBORD,
         "Commodity Rail", WHITE, 15, "present everywhere, earns nothing", GREY)
    quad(s, xRcol, yBrow, qw, qh, RGBColor(0x0A, 0x1A, 0x0A), RED,
         "WhatsApp Pay Today", RED, 16, "< 1% share", WHITE)

    # red "you are here" dot in bottom-right (below the top-anchored label)
    d = oval(s, xRcol + qw // 2, yBrow + qh - Inches(0.5), Inches(0.1), RED)
    glow_shape(d, hx(RED), rad=63500, alpha=55000)

    # diagonal path arrow: bottom-left corner → top-right corner (passes through empty lower halves)
    arrow(s, xL + Inches(0.35), yB - Inches(0.28), xR - Inches(0.35), yT + Inches(0.34),
          GREEN, Pt(3.5), big=True, glow=True)
    # label chip at the empty centre cross-gap (drawn last, covers the arrow crossing)
    cxm, cym = (xL + xR) // 2, (yT + yB) // 2
    cw2, chh = Inches(3.8), Inches(0.52)
    chip = rrect(s, cxm - cw2 // 2, cym - chh // 2, cw2, chh, BG, GREEN, Pt(1.25), radius=0.45)
    ctf = chip.text_frame; ctf.word_wrap = True; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = ctf.margin_right = Pt(6); ctf.margin_top = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = "The path: CTWA closed-loop checkout + merchant lending"
    cr.font.size, cr.font.bold, cr.font.name = Pt(10), True, "Calibri"; cr.font.color.rgb = GREEN

    # axis labels
    xax = s.shapes.add_textbox(xL, yB + Inches(0.14), xR - xL, Inches(0.3))
    xtf = xax.text_frame; xtf.word_wrap = False; xtf.margin_top = xtf.margin_bottom = 0
    xp = xtf.paragraphs[0]; xp.alignment = PP_ALIGN.CENTER
    xr = xp.add_run(); xr.text = "META'S STRATEGIC COMMITMENT          LOW  →  HIGH"
    xr.font.size, xr.font.bold, xr.font.name = Pt(10.5), True, "Calibri"; xr.font.color.rgb = GREY_DK
    yax = s.shapes.add_textbox(Inches(-1.6), Inches(3.62), Inches(4.6), Inches(0.34))
    yax.rotation = 270
    ytf = yax.text_frame; ytf.word_wrap = False; ytf.margin_top = ytf.margin_bottom = 0
    yp = ytf.paragraphs[0]; yp.alignment = PP_ALIGN.CENTER
    yr = yp.add_run(); yr.text = "MONETIZATION        ZERO-MDR ↓   /   LTV-POSITIVE ↑"
    yr.font.size, yr.font.bold, yr.font.name = Pt(10.5), True, "Calibri"; yr.font.color.rgb = GREY_DK

    footer(s, 6)

# ═════════════════════════════════════════════════════════
# SLIDE 12 — SOLUTIONS VALIDATED (RICE)
# ═════════════════════════════════════════════════════════
def slide12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "RIGOR")
    title(s, "Solutions Prioritized — RICE × Stress-Test", size=25)
    textbox(s, "RICE = (Reach × Impact × Confidence) ÷ Effort    ·    R / I / C scored 0–10,  Effort 1–5 (lower = easier to ship)",
            Inches(0.4), Inches(0.96), Inches(12.5), Inches(0.3), size=10.5, color=GREY_DK, italic=True)

    WEDGE = RGBColor(0x00, 0xE6, 0x76); HABIT = RGBColor(0x00, 0xB8, 0x94); MOAT = RGBColor(0xFF, 0xB3, 0x00)
    rows = [
        ("5th Tab + Camera Intercept", "WEDGE · ship now", WEDGE, 10, 7, 9, 1, 630,
         "Zero-CAC fix to the #1 cause of counter-side abandonment — reaches every user instantly.",
         "NPCI may restrict default-payment UI patterns."),
        ("CTWA Closed-Loop Checkout", "WEDGE · ship now", WEDGE, 9, 10, 8, 2, 360,
         "Turns a ₹0 UPI txn into an attributed ad-sale — adds LTV without changing UPI economics.",
         "Weak checkout UX pushes CTWA ad spend to rivals."),
        ("JioMart Anchor P2M", "HABIT · frequency", HABIT, 8, 9, 7, 3, 168,
         "One mega-partner manufactures 40–160M P2M txns/mo — builds the daily-default habit.",
         "Reliance dependency; Jio could exit the deal."),
        ("Hongbao Viral Loop", "HABIT · frequency", HABIT, 9, 6, 5, 2, 135,
         "Social money-gifting drives P2P virality at near-zero acquisition cost.",
         "Gifting behaviour unproven in India; novelty may fade."),
        ("Merchant Lending Engine", "MOAT · strategic", MOAT, 6, 10, 7, 4, 105,
         "Highest LTV per merchant; a compounding credit moat rivals can't match in-chat.",
         "NBFC licensing, credit-risk exposure, slow to stand up."),
        ("Inbound Remittance (NRI)", "MOAT · strategic", MOAT, 5, 10, 6, 4, 75,
         "Only player with a 180-country graph; greenfield slice of a $135B corridor.",
         "FEMA / RBI cross-border rules + FX-compliance load."),
    ]

    tx = Inches(0.35)
    colw = [Inches(2.5), Inches(0.44), Inches(0.44), Inches(0.44), Inches(0.44),
            Inches(0.95), Inches(3.7), Inches(3.7)]
    heads = ["Solution", "R", "I", "C", "E", "RICE", "Why it works", "Kill risk"]
    xcols = []; cxp = tx
    for w in colw: xcols.append(cxp); cxp += w
    total_w = sum(colw)

    hy, hh = Inches(1.34), Inches(0.46)
    rect(s, tx, hy, total_w, hh, RGBColor(0x1B, 0x5E, 0x20))
    for i, hd in enumerate(heads):
        al = PP_ALIGN.LEFT if i in (0, 6, 7) else PP_ALIGN.CENTER
        pad = Inches(0.16) if i in (0, 6, 7) else Inches(0)
        textbox(s, hd, xcols[i] + pad, hy, colw[i] - pad, hh, size=10.5, color=WHITE,
                bold=True, align=al, anchor=MSO_ANCHOR.MIDDLE)

    ry = hy + hh; rh = Inches(0.76)
    for ri, (name, role, rc, R, I, C, E, rice, why, kill) in enumerate(rows):
        rect(s, tx, ry, total_w, rh, RGBColor(0x12, 0x18, 0x1F) if ri % 2 == 0 else RGBColor(0x0C, 0x10, 0x15))
        rect(s, tx, ry, Emu(52000), rh, rc)  # role accent bar
        sb = s.shapes.add_textbox(xcols[0] + Inches(0.18), ry, colw[0] - Inches(0.22), rh)
        stf = sb.text_frame; stf.word_wrap = True; stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        stf.margin_left = stf.margin_right = stf.margin_top = stf.margin_bottom = 0
        para(stf, name, 11, WHITE, bold=True, first=True)
        para(stf, role, 8, rc, bold=True, before=1)
        for j, val in enumerate([R, I, C, E]):
            textbox(s, str(val), xcols[1 + j], ry, colw[1 + j], rh, size=11.5, color=WHITE,
                    bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cw_ = colw[5] - Inches(0.18); chp = Inches(0.42)
        cxx = xcols[5] + Inches(0.09); cyy = ry + (rh - chp) // 2
        rrect(s, cxx, cyy, cw_, chp, rc, None, radius=0.28)
        textbox(s, str(rice), cxx, cyy, cw_, chp, size=13, color=RGBColor(0x06, 0x12, 0x0A),
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        wb = s.shapes.add_textbox(xcols[6] + Inches(0.1), ry, colw[6] - Inches(0.2), rh)
        wtf = wb.text_frame; wtf.word_wrap = True; wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        wtf.margin_left = wtf.margin_right = wtf.margin_top = wtf.margin_bottom = 0
        para(wtf, why, 9.5, RGBColor(0xA5, 0xD6, 0xA7), first=True)
        kb = s.shapes.add_textbox(xcols[7] + Inches(0.1), ry, colw[7] - Inches(0.2), rh)
        ktf = kb.text_frame; ktf.word_wrap = True; ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ktf.margin_left = ktf.margin_right = ktf.margin_top = ktf.margin_bottom = 0
        para(ktf, kill, 9.5, RGBColor(0xE0, 0xA0, 0xA0), first=True)
        ry += rh

    iy = ry + Inches(0.1)
    rrect(s, tx, iy, total_w, Inches(0.52), RGBColor(0x10, 0x16, 0x12), GREEN, Pt(1), radius=0.12)
    ib = s.shapes.add_textbox(tx + Inches(0.25), iy, total_w - Inches(0.5), Inches(0.52))
    itf = ib.text_frame; itf.word_wrap = True; itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    itf.margin_left = itf.margin_right = itf.margin_top = itf.margin_bottom = 0
    pp = itf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    def rn(t, c, b=True, sz=10.5):
        r = pp.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
        r.font.name = "Calibri"; r.font.color.rgb = c
    rn("THE READ:   ", GREY)
    rn("Wedge", WEDGE); rn(" ships first — cheap & universal.     ", WHITE, False)
    rn("Habit", HABIT); rn(" manufactures frequency.     ", WHITE, False)
    rn("Moat", MOAT); rn(" carries the highest impact, built after the habit exists.", WHITE, False)
    footer(s, 12)

# ═════════════════════════════════════════════════════════
# SLIDE 13 — HOW ₹0 BECOMES A BUSINESS (LTV column chart)
# ═════════════════════════════════════════════════════════
def slide13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    label_pill(s, "STRATEGY")
    title(s, "How ₹0 Becomes a Business", size=27)
    textbox(s, "LTV per active user / year (₹) — and exactly where each rupee comes from",
            Inches(0.4), Inches(0.96), Inches(12), Inches(0.3), size=11, color=GREY_DK, italic=True)

    AD = RGBColor(0x00, 0xE6, 0x76); LEND = RGBColor(0x00, 0x89, 0x7B)
    INS = RGBColor(0x4C, 0xAF, 0x50); FX = RGBColor(0xFF, 0xB3, 0x00)
    streams = [
        ("Ad-attributed commerce (CTWA)", AD,   [0, 120, 400], "Meta's take-rate on CTWA-driven sales"),
        ("Merchant lending",              LEND, [0, 0, 500],   "interest on working-capital credit"),
        ("Insurance & credit",            INS,  [0, 0, 180],   "commission on attached products"),
        ("Inbound remittance (FX)",       FX,   [0, 0, 170],   "FX spread on NRI inflows"),
    ]
    cats = ["TODAY", "12–18 MONTHS", "24–36 MONTHS"]
    totals = [0, 120, 1250]
    tlabels = ["₹0", "₹50–200", "₹500–2,000"]
    tcolors = [RGBColor(0xFF, 0x6B, 0x6B), WHITE, GREEN]
    vmax = 1400

    Cx, Cy, Cw, Ch = Inches(0.35), Inches(1.55), Inches(7.85), Inches(4.25)
    add_stacked_column(s, Cx, Cy, Cw, Ch, cats,
                       [(nm, vals) for nm, col, vals, mech in streams],
                       [col for nm, col, vals, mech in streams], vmax=vmax, gap=85)

    for i in range(3):
        cxc = Cx + Cw * (2 * i + 1) // 6
        topy = Cy + int(Ch * (1 - totals[i] / vmax))
        ly = (Cy + Ch - Inches(0.55)) if i == 0 else (topy - Inches(0.5))
        textbox(s, tlabels[i], cxc - Inches(1.45), ly, Inches(2.9), Inches(0.42),
                size=17, color=tcolors[i], bold=True, align=PP_ALIGN.CENTER,
                glow=(("00E676", 50000, 50000) if i == 2 else None))
        textbox(s, cats[i], cxc - Inches(1.45), Cy + Ch + Inches(0.1), Inches(2.9), Inches(0.32),
                size=11.5, color=GREY, bold=True, align=PP_ALIGN.CENTER)

    # right panel — revenue stack / inference (doubles as the chart legend)
    rx, ry0, rw = Inches(8.45), Inches(1.55), Inches(4.5)
    rrect(s, rx, ry0, rw, Inches(4.3), RGBColor(0x0D, 0x10, 0x0D), RGBColor(0x1E, 0x2D, 0x24), Pt(1), radius=0.05)
    hb = s.shapes.add_textbox(rx + Inches(0.28), ry0 + Inches(0.2), rw - Inches(0.56), Inches(0.8))
    htf = hb.text_frame; htf.word_wrap = True; htf.margin_left = htf.margin_right = htf.margin_top = htf.margin_bottom = 0
    para(htf, "WHERE THE ₹ COMES FROM", 13, GREEN, bold=True, first=True)
    para(htf, "Zero-MDR removed the toll on the rail. The rail was never where the money was.",
         10, GREY, italic=True, before=4)
    sy = ry0 + Inches(1.2)
    for nm, col, vals, mech in streams:
        oval(s, rx + Inches(0.4), sy + Inches(0.18), Inches(0.08), col)
        lb = s.shapes.add_textbox(rx + Inches(0.64), sy, rw - Inches(0.85), Inches(0.56))
        ltf = lb.text_frame; ltf.word_wrap = True; ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ltf.margin_left = ltf.margin_right = ltf.margin_top = ltf.margin_bottom = 0
        para(ltf, nm, 10.5, WHITE, bold=True, first=True)
        para(ltf, mech, 9, GREY, before=1)
        sy += Inches(0.63)
    st = rrect(s, rx + Inches(0.3), sy + Inches(0.04), rw - Inches(0.6), Inches(0.5),
               RGBColor(0x0A, 0x22, 0x0A), GREEN, Pt(1.25), radius=0.28)
    stf2 = st.text_frame; stf2.vertical_anchor = MSO_ANCHOR.MIDDLE; stf2.word_wrap = True
    stf2.margin_left = stf2.margin_right = stf2.margin_top = stf2.margin_bottom = 0
    para(stf2, "✓  EVERY STREAM IS MDR-INDEPENDENT", 11, GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)

    banner(s, "₹0 on the rail — but ads, lending, insurance and FX are all MDR-independent. That's the LTV no competitor can build inside an ad platform.",
           Inches(6.45))
    footer(s, 13)

# ═════════════════════════════════════════════════════════
# SLIDE 14 — META MOVED (epilogue)
# ═════════════════════════════════════════════════════════
def slide14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    SW = prs.slide_width
    label_pill(s, "EPILOGUE")
    textbox(s, "Meta Moved.", Inches(0), Inches(0.34), SW, Inches(0.8), size=40, color=GREEN,
            bold=True, align=PP_ALIGN.CENTER, glow=("00E676", 88900, 60000))
    rect(s, (SW - Inches(2.2)) // 2, Inches(1.2), Inches(2.2), Emu(28000), GREEN)
    textbox(s, "Three days after this analysis was written:", Inches(0.5), Inches(1.42),
            Inches(8), Inches(0.35), size=14, color=GREY, italic=True)

    lx, lw = Inches(0.5), Inches(7.85)
    # CRED investment card
    rrect(s, lx, Inches(1.92), lw, Inches(1.15), CARD, GREEN, Pt(1.75), radius=0.06)
    c1 = s.shapes.add_textbox(lx + Inches(0.3), Inches(1.92), lw - Inches(0.6), Inches(1.15))
    t1 = c1.text_frame; t1.word_wrap = True; t1.vertical_anchor = MSO_ANCHOR.MIDDLE
    t1.margin_left = t1.margin_right = t1.margin_top = t1.margin_bottom = 0
    r = para(t1, "₹7,500 Cr", 34, GREEN, bold=True, first=True); glow_run(r, "00E676", rad=63500, alpha=50000)
    para(t1, "Meta-led investment in CRED ($900M) — Jun 22, 2026", 12, WHITE, before=2)
    # Kunal Shah card
    rrect(s, lx, Inches(3.22), lw, Inches(1.62), CARD, GREEN, Pt(1.75), radius=0.06)
    c2 = s.shapes.add_textbox(lx + Inches(0.3), Inches(3.22), lw - Inches(0.6), Inches(1.62))
    t2 = c2.text_frame; t2.word_wrap = True; t2.vertical_anchor = MSO_ANCHOR.MIDDLE
    t2.margin_left = t2.margin_right = t2.margin_top = t2.margin_bottom = 0
    para(t2, "Kunal Shah", 27, WHITE, bold=True, first=True)
    para(t2, "CRED founder appointed Global Head of WhatsApp, succeeding Will Cathcart", 12, WHITE, before=3)
    para(t2, "•  First Indian CEO of WhatsApp globally", 11, GREY, before=4)
    para(t2, "•  Background: FreeCharge → CRED → WhatsApp Pay", 11, GREY, before=2)
    # callout
    coy = Inches(5.02)
    rrect(s, lx, coy, lw, Inches(1.08), RGBColor(0x0A, 0x2A, 0x0A), None, radius=0.06)
    rect(s, lx, coy, Emu(48000), Inches(1.08), GREEN)
    co = s.shapes.add_textbox(lx + Inches(0.25), coy, lw - Inches(0.45), Inches(1.08))
    tco = co.text_frame; tco.word_wrap = True; tco.vertical_anchor = MSO_ANCHOR.MIDDLE
    tco.margin_left = tco.margin_right = tco.margin_top = tco.margin_bottom = 0
    para(tco, "His stated challenge: making WhatsApp Pay matter in India.", 12.5, WHITE, bold=True, first=True)
    para(tco, "Our analysis said why it failed. Meta just hired someone to fix exactly that.", 12, GREEN, italic=True, before=3)

    # right — thesis validated
    rx, rw = Inches(8.6), Inches(4.33)
    rrect(s, rx, Inches(1.92), rw, Inches(4.18), RGBColor(0x11, 0x11, 0x11), None, radius=0.04)
    rect(s, rx, Inches(1.92), rw, Emu(42000), GREEN)
    th = s.shapes.add_textbox(rx + Inches(0.28), Inches(2.05), rw - Inches(0.56), Inches(3.95))
    tt = th.text_frame; tt.word_wrap = True; tt.vertical_anchor = MSO_ANCHOR.MIDDLE
    tt.margin_left = tt.margin_right = tt.margin_top = tt.margin_bottom = 0
    para(tt, "THE THESIS, VALIDATED", 14, GREEN, bold=True, first=True)
    for b in ["Root cause: Meta's will to invest — not UPI economics",
              "Strategic pivot: Commerce OS via CTWA",
              "Global moat: WhatsApp's unique cross-border graph"]:
        p = tt.add_paragraph(); p.space_before = Pt(13)
        rk = p.add_run(); rk.text = "✓  "; rk.font.size = Pt(12.5); rk.font.bold = True
        rk.font.color.rgb = GREEN; rk.font.name = "Calibri"
        rb = p.add_run(); rb.text = b; rb.font.size = Pt(12.5); rb.font.color.rgb = WHITE; rb.font.name = "Calibri"
    p = tt.add_paragraph(); p.space_before = Pt(20)
    rr = p.add_run(); rr.text = "Every variable was downstream of one decision."
    rr.font.size, rr.font.italic, rr.font.name = Pt(11.5), True, "Calibri"; rr.font.color.rgb = GREY

    # sources strip
    rect(s, Inches(0), Inches(7.06), SW, Inches(0.44), RGBColor(0x0E, 0x1A, 0x0E))
    textbox(s, "Sources: NPCI · RBI payment-system reports · PhonePe FY25 filings · Meta earnings calls · TechCrunch, Bloomberg, CNBC — Jun 22, 2026",
            Inches(0.4), Inches(7.06), Inches(12.5), Inches(0.44), size=9, color=GREY, anchor=MSO_ANCHOR.MIDDLE)

# ═════════════════════════════════════════════════════════
# SLIDES 7–11 — IMAGE SLIDES (prototypes / synthesis)
# ═════════════════════════════════════════════════════════
MEDIA = r"c:\Users\91960\Claude\Projects\PM journey\whatsapp case study\deck_media"
IMG = {
    "synthesis":  MEDIA + r"\synthesis.png",
    "ctwa":       MEDIA + r"\sol1_ctwa.png",
    "jiomart":    MEDIA + r"\sol2_jiomart.png",
    "fifthtab":   MEDIA + r"\sol3_5thtab.png",
    "remittance": MEDIA + r"\sol4_remittance.png",
}

def image_slide(prs, img_path, pill_text, n):
    """Full-bleed self-contained image, aspect preserved, with pill + footer on the dark margins."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    img_h = Inches(6.4)
    img_w = Inches(6.4 * 1376 / 768)          # preserve 1376x768 aspect → 11.47"
    img_x = (prs.slide_width - img_w) // 2
    img_y = Inches(0.62)
    pic = s.shapes.add_picture(img_path, img_x, img_y, img_w, img_h)
    pic.line.color.rgb = GREEN_MID; pic.line.width = Pt(1.25)   # subtle frame
    label_pill(s, pill_text)
    footer(s, n)
    return s

# ── BUILD ────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide1(prs); slide2(prs); slide3(prs); slide4(prs); slide5(prs); slide6(prs)
    image_slide(prs, IMG["synthesis"],  "STRATEGY",   7)
    image_slide(prs, IMG["ctwa"],       "SOLUTION 01", 8)
    image_slide(prs, IMG["jiomart"],    "SOLUTION 02", 9)
    image_slide(prs, IMG["fifthtab"],   "SOLUTION 03", 10)
    image_slide(prs, IMG["remittance"], "SOLUTION 04", 11)
    slide12(prs); slide13(prs); slide14(prs)
    out = r"c:\Users\91960\Claude\Projects\PM journey\whatsapp case study\WhatsApp_Pay_Final.pptx"
    prs.save(out); print("Saved:", out)

if __name__ == "__main__":
    build()
