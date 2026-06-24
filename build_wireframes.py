# -*- coding: utf-8 -*-
"""
CTWA Closed-Loop Checkout — UX Wireframes (phone-frame flow), python-pptx only.
Usage: python build_wireframes.py
Output: WhatsApp_Pay_CTWA_Wireframes.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


def C(h): return RGBColor.from_string(h)
PRIMARY = C('075E54'); TEAL = C('128C7E'); ACCENT = C('25D366')
DARK = C('0B141A'); LIGHT = C('F0F2F5'); WHITE = C('FFFFFF')
BORDER = C('E0E0E0'); RED = C('E63946'); AMBER = C('F4A261'); GREY = C('8D99AE')
INK = C('111B21'); SUBTLE = C('667781')
CHAT_BG = C('ECE5DD'); OUT = C('DCF8C6'); PHONE = C('0B141A')
HILITE = C('E8F7F0')


def R(t, s, b, c, i=False): return (t, s, b, c, i)
def P(runs, align=None, sa=None, sb=None, ln=None):
    d = {'runs': runs}
    if align is not None: d['align'] = align
    if sa is not None: d['sa'] = sa
    if sb is not None: d['sb'] = sb
    if ln is not None: d['ln'] = ln
    return d


def tx(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align)
        if 'sa' in para: p.space_after = Pt(para['sa'])
        if 'sb' in para: p.space_before = Pt(para['sb'])
        if 'ln' in para: p.line_spacing = para['ln']
        for run in para['runs']:
            r = p.add_run(); r.text = run[0]
            r.font.size = Pt(run[1]); r.font.bold = run[2]
            r.font.color.rgb = run[3]; r.font.name = 'Calibri'
            if len(run) > 4 and run[4]: r.font.italic = True
    return box


def rr(s, x, y, w, h, fill, line=None, radius=0.12, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    return shp


def rc(s, x, y, w, h, fill, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def ov(s, x, y, d, fill, line=None, lw=1.0, dh=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(dh or d))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def glyph(s, x, y, d, ch, color, size, fill=None):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if fill is None: o.fill.background()
    else: o.fill.solid(); o.fill.fore_color.rgb = fill
    o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ch; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color; r.font.name = 'Calibri'
    return o


def pill(s, x, y, w, h, text, fill=PRIMARY, tcol=WHITE, size=8, line=None):
    shp = rr(s, x, y, w, h, fill, line=line, radius=0.5)
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = tcol; r.font.name = 'Calibri'
    return shp


def pin(s, cx, cy, num, color=ACCENT, d=0.26):
    glyph(s, cx - d / 2, cy - d / 2, d, str(num), WHITE, 11, fill=color)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    o.fill.background(); o.line.color.rgb = WHITE; o.line.width = Pt(1.5); o.shadow.inherit = False


def arrow(s, x1, y1, x2, y2, color, weight=2.0):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    ln.append(parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>' % nsdecls('a')))
    conn.shadow.inherit = False
    return conn


def bg(s, color):
    f = s.background.fill; f.solid(); f.fore_color.rgb = color


def header(s, title, kicker):
    rc(s, 0, 0, 13.333, 1.0, PRIMARY)
    rc(s, 0, 1.0, 13.333, 0.06, ACCENT)
    tx(s, 0.5, 0.1, 9.7, 0.82, [P([R(title, 24, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
    p = rr(s, 10.85, 0.3, 2.0, 0.42, TEAL, radius=0.5)
    tf = p.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    r = pr.add_run(); r.text = kicker; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = 'Calibri'


def footer(s, n, total=6):
    tx(s, 0.5, 7.08, 9.0, 0.34, [P([R('WhatsApp Pay × CTWA — Closed-Loop Checkout · UX Wireframes', 9, False, GREY)])],
       anchor=MSO_ANCHOR.MIDDLE)
    tx(s, 10.8, 7.08, 2.03, 0.34, [P([R('%d / %d' % (n, total), 10, True, GREY)], align=PP_ALIGN.RIGHT)],
       anchor=MSO_ANCHOR.MIDDLE)


# ---------- phone frame ----------
def phone(s, x, y, w=1.95, h=3.75):
    rr(s, x, y, w, h, PHONE, radius=0.09)
    mx, mt, mb = 0.055, 0.07, 0.07
    sx, sy, sw, sh = x + mx, y + mt, w - 2 * mx, h - mt - mb
    rr(s, sx, sy, sw, sh, WHITE, radius=0.04)
    return sx, sy, sw, sh


def wa_header(s, sx, sy, sw, name='Bharat Kirana', sub='business account'):
    rc(s, sx, sy, sw, 0.34, PRIMARY)
    tx(s, sx + 0.02, sy, 0.16, 0.34, [P([R('‹', 13, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
    ov(s, sx + 0.2, sy + 0.05, 0.24, C('B8D8CE'))
    tx(s, sx + 0.48, sy + 0.02, sw - 0.7, 0.34, [
        P([R(name, 7.5, True, WHITE)], sa=0, ln=0.9),
        P([R(sub, 5.5, False, C('CDEDE4'))], ln=0.9)], anchor=MSO_ANCHOR.MIDDLE)


def input_bar(s, sx, sy, sw, sh):
    rr(s, sx + 0.05, sy + sh - 0.3, sw - 0.42, 0.24, C('F5F5F5'), line=BORDER, radius=0.5)
    tx(s, sx + 0.12, sy + sh - 0.31, sw - 0.55, 0.24, [P([R('Message', 6.5, False, SUBTLE)])], anchor=MSO_ANCHOR.MIDDLE)
    ov(s, sx + sw - 0.33, sy + sh - 0.31, 0.26, ACCENT)
    tx(s, sx + sw - 0.33, sy + sh - 0.31, 0.26, 0.26, [P([R('➤', 8, True, WHITE)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)


# ---------- screen drawers ----------
def draw_ad(s, sx, sy, sw, sh):
    # status + app bar
    rc(s, sx, sy, sw, 0.26, WHITE)
    tx(s, sx + 0.06, sy, 0.5, 0.26, [P([R('9:41', 6, True, INK)])], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx, sy, sw - 0.06, 0.26, [P([R('◉ Instagram', 7.5, True, INK)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    rc(s, sx, sy + 0.26, sw, 0.012, BORDER)
    # post header
    ov(s, sx + 0.07, sy + 0.34, 0.24, C('CFE9DF'))
    tx(s, sx + 0.36, sy + 0.32, sw - 0.4, 0.3, [
        P([R('bharat_kirana', 7.5, True, INK)], ln=0.9),
        P([R('Sponsored', 6, False, SUBTLE)], ln=0.9)])
    # image
    rc(s, sx + 0.05, sy + 0.66, sw - 0.1, 1.55, C('E4E7EB'))
    tx(s, sx + 0.05, sy + 0.66, sw - 0.1, 1.55, [
        P([R('▦', 16, False, C('B7BEC7'))], align=PP_ALIGN.CENTER, ln=0.9),
        P([R('product image', 6.5, False, C('9AA3AD'))], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    # actions
    tx(s, sx + 0.07, sy + 2.24, sw - 0.14, 0.2, [P([R('♡    ◯    ➤', 9, False, INK)])], anchor=MSO_ANCHOR.MIDDLE)
    # CTA
    pill(s, sx + 0.05, sy + 2.5, sw - 0.1, 0.34, '🟢  Shop on WhatsApp', ACCENT, WHITE, 7.5)
    # caption
    tx(s, sx + 0.07, sy + 2.92, sw - 0.14, 0.6, [
        P([R('Aashirvaad Atta 10kg — ₹545.', 6.8, True, INK)], ln=0.95),
        P([R('Tap to order on WhatsApp ↗', 6.5, False, SUBTLE)], ln=0.95)])


def draw_chat_catalog(s, sx, sy, sw, sh):
    wa_header(s, sx, sy, sw)
    rc(s, sx, sy + 0.34, sw, sh - 0.34, CHAT_BG)
    # greeting bubble
    rr(s, sx + 0.07, sy + 0.46, 1.45, 0.4, WHITE, radius=0.18)
    tx(s, sx + 0.12, sy + 0.47, 1.36, 0.4, [P([R('Namaste! 👋 Here you go:', 6.8, False, INK)], ln=0.95)],
       anchor=MSO_ANCHOR.MIDDLE)
    # catalog card
    cardx, cardy, cardw = sx + 0.07, sy + 0.94, 1.6
    rr(s, cardx, cardy, cardw, 1.85, WHITE, line=BORDER, radius=0.1)
    rc(s, cardx + 0.07, cardy + 0.08, cardw - 0.14, 0.78, C('E4E7EB'))
    tx(s, cardx + 0.07, cardy + 0.08, cardw - 0.14, 0.78, [P([R('▦ image', 6.5, False, C('9AA3AD'))], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    tx(s, cardx + 0.08, cardy + 0.9, cardw - 0.16, 0.5, [
        P([R('Aashirvaad Atta 10kg', 7.5, True, INK)], ln=0.95),
        P([R('₹545', 9, True, PRIMARY)], ln=1.0)])
    pill(s, cardx + 0.08, cardy + 1.45, cardw - 0.16, 0.3, '＋  Add to cart', PRIMARY, WHITE, 7)
    input_bar(s, sx, sy, sw, sh)


def draw_cart(s, sx, sy, sw, sh):
    wa_header(s, sx, sy, sw)
    rc(s, sx, sy + 0.34, sw, sh - 0.34, CHAT_BG)
    cardx, cardy, cardw = sx + 0.1, sy + 0.5, 1.6
    rr(s, cardx, cardy, cardw, 2.35, WHITE, line=BORDER, radius=0.08)
    tx(s, cardx + 0.1, cardy + 0.07, cardw - 0.2, 0.28, [P([R('🛒  Your cart', 8, True, INK)])], anchor=MSO_ANCHOR.MIDDLE)
    rc(s, cardx + 0.1, cardy + 0.36, cardw - 0.2, 0.01, BORDER)
    items = [('Aashirvaad Atta 10kg ×1', '₹545'), ('Tata Salt 1kg ×2', '₹56')]
    iy = cardy + 0.44
    for nm, pr in items:
        tx(s, cardx + 0.1, iy, cardw - 0.55, 0.3, [P([R(nm, 6.8, False, INK)], ln=0.95)], anchor=MSO_ANCHOR.MIDDLE)
        tx(s, cardx + cardw - 0.5, iy, 0.42, 0.3, [P([R(pr, 6.8, True, INK)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        iy += 0.34
    rc(s, cardx + 0.1, iy + 0.02, cardw - 0.2, 0.01, BORDER)
    tx(s, cardx + 0.1, iy + 0.08, cardw - 0.5, 0.3, [P([R('Subtotal', 7.2, True, INK)])], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, cardx + cardw - 0.5, iy + 0.08, 0.42, 0.3, [P([R('₹601', 8, True, PRIMARY)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    pill(s, cardx + 0.1, iy + 0.46, cardw - 0.2, 0.32, 'Checkout  →', PRIMARY, WHITE, 8)
    input_bar(s, sx, sy, sw, sh)


def draw_checkout(s, sx, sy, sw, sh):
    wa_header(s, sx, sy, sw)
    rc(s, sx, sy + 0.34, sw, sh - 0.34, C('CFCABF'))   # dimmed chat behind
    # bottom sheet
    shy = sy + 0.7
    rr(s, sx, shy, sw, sh - (shy - sy), WHITE, radius=0.06)
    ov(s, sx + sw / 2 - 0.12, shy + 0.05, 0.24, C('D6D6D6'), dh=0.04)
    tx(s, sx + 0.1, shy + 0.12, sw - 0.2, 0.26, [P([R('Checkout', 9, True, INK)])])
    tx(s, sx + 0.1, shy + 0.12, sw - 0.2, 0.26, [P([R('2 items · ₹601', 6.5, False, SUBTLE)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    # address
    rr(s, sx + 0.1, shy + 0.42, sw - 0.2, 0.34, C('F7F7F7'), line=BORDER, radius=0.2)
    tx(s, sx + 0.16, shy + 0.42, sw - 0.55, 0.34, [P([R('📍 Home — Flat 4B, Indira Nagar', 6.3, False, INK)], ln=0.9)], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx + sw - 0.5, shy + 0.42, 0.4, 0.34, [P([R('Change', 6, True, PRIMARY)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    # payment label
    tx(s, sx + 0.12, shy + 0.82, sw - 0.2, 0.2, [P([R('PAYMENT METHOD', 6, True, SUBTLE)])])
    # WhatsApp Pay (default / selected) — highlighted
    rr(s, sx + 0.1, shy + 1.02, sw - 0.2, 0.42, HILITE, line=ACCENT, radius=0.16, lw=1.5)
    ov(s, sx + 0.17, shy + 1.14, 0.16, WHITE, line=ACCENT, lw=1.5)
    ov(s, sx + 0.195, shy + 1.165, 0.11, ACCENT)
    tx(s, sx + 0.4, shy + 1.02, sw - 0.5, 0.42, [
        P([R('WhatsApp Pay', 7.5, True, INK)], ln=0.9),
        P([R('UPI · default', 5.8, False, PRIMARY)], ln=0.9)], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx + sw - 0.32, shy + 1.02, 0.22, 0.42, [P([R('✓', 9, True, ACCENT)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    # other UPI apps (unselected)
    rr(s, sx + 0.1, shy + 1.5, sw - 0.2, 0.36, WHITE, line=BORDER, radius=0.16)
    ov(s, sx + 0.17, shy + 1.6, 0.16, WHITE, line=GREY, lw=1.2)
    tx(s, sx + 0.4, shy + 1.5, sw - 0.5, 0.36, [P([R('Other UPI apps', 7, False, SUBTLE)])], anchor=MSO_ANCHOR.MIDDLE)
    # total + pay
    tx(s, sx + 0.12, shy + 1.96, sw - 0.2, 0.24, [P([R('To pay', 7, False, INK)])], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx + 0.12, shy + 1.96, sw - 0.22, 0.24, [P([R('₹601', 9, True, INK)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    pill(s, sx + 0.1, shy + 2.24, sw - 0.2, 0.34, 'Pay ₹601', ACCENT, WHITE, 8.5)


def draw_pin(s, sx, sy, sw, sh):
    rc(s, sx, sy, sw, sh, C('F4F6F8'))
    rc(s, sx, sy, sw, 0.5, PRIMARY)
    tx(s, sx, sy + 0.04, sw, 0.46, [
        P([R('🔒 Enter UPI PIN', 8, True, WHITE)], align=PP_ALIGN.CENTER, ln=0.9),
        P([R('HDFC Bank ····1234', 6, False, C('CDEDE4'))], align=PP_ALIGN.CENTER, ln=0.9)], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx + 0.1, sy + 0.6, sw - 0.2, 0.34, [
        P([R('Paying ₹601 to Bharat Kirana', 6.8, True, INK)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    # PIN dots
    dotx = sx + sw / 2 - 0.45
    for i in range(6):
        ov(s, dotx + i * 0.16, sy + 1.0, 0.1, INK if i < 4 else None, line=(None if i < 4 else GREY), lw=1.0)
    # keypad
    keys = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '', '0', '⌫']
    kw, kh, gx, gy = 0.42, 0.3, 0.06, 0.06
    gridw = 3 * kw + 2 * gx
    startx = sx + (sw - gridw) / 2
    starty = sy + 1.35
    for idx, k in enumerate(keys):
        col = idx % 3; row = idx // 3
        kx = startx + col * (kw + gx); ky = starty + row * (kh + gy)
        if k == '': continue
        rr(s, kx, ky, kw, kh, WHITE, line=BORDER, radius=0.12)
        tx(s, kx, ky, kw, kh, [P([R(k, 9, True, INK)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tx(s, sx, sy + sh - 0.22, sw, 0.2, [P([R('Secured by UPI · NPCI', 5.8, False, SUBTLE)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)


def draw_success(s, sx, sy, sw, sh):
    wa_header(s, sx, sy, sw)
    rc(s, sx, sy + 0.34, sw, sh - 0.34, CHAT_BG)
    # success chip (centered system)
    glyph(s, sx + sw / 2 - 0.22, sy + 0.5, 0.44, '✓', WHITE, 16, fill=ACCENT)
    tx(s, sx + 0.1, sy + 1.0, sw - 0.2, 0.24, [P([R('Payment Successful', 8, True, PRIMARY)], align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)
    # receipt bubble
    cardx, cardy, cardw = sx + 0.12, sy + 1.34, 1.55
    rr(s, cardx, cardy, cardw, 1.62, WHITE, line=BORDER, radius=0.1)
    tx(s, cardx + 0.1, cardy + 0.08, cardw - 0.2, 1.5, [
        P([R('Order #BK-10293  ✓', 7.5, True, INK)], sa=2, ln=0.95),
        P([R('Aashirvaad Atta 10kg ×1', 6.5, False, INK)], ln=0.95),
        P([R('Tata Salt 1kg ×2', 6.5, False, INK)], sa=3, ln=0.95),
        P([R('Paid ₹601 · WhatsApp Pay', 6.8, True, PRIMARY)], sa=4, ln=0.95)])
    pill(s, cardx + 0.1, cardy + 1.2, cardw - 0.2, 0.3, 'Track order  →', PRIMARY, WHITE, 7)
    input_bar(s, sx, sy, sw, sh)


# ---------- caption under a phone ----------
def phone_caption(s, cx, num, title, body, color=ACCENT):
    glyph(s, cx - 0.14, 5.12, 0.28, str(num), WHITE, 12, fill=color)
    tx(s, cx - 1.85, 5.5, 3.7, 1.3, [
        P([R(title, 11, True, INK)], align=PP_ALIGN.CENTER, sa=3, ln=1.0),
        P([R(body, 9.5, False, SUBTLE)], align=PP_ALIGN.CENTER, ln=1.05)], align=PP_ALIGN.CENTER)


# =====================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # -------- SLIDE 1 : TITLE --------
    s = prs.slides.add_slide(blank); bg(s, DARK)
    tx(s, 1.0, 1.7, 11.33, 2.0, [
        P([R('CTWA Closed-Loop Checkout', 38, True, WHITE)], align=PP_ALIGN.CENTER, sa=6),
        P([R('UX Wireframes & End-to-End User Flow', 20, True, ACCENT)], align=PP_ALIGN.CENTER)],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, 1.5, 3.7, 10.33, 0.5, [P([R(
        'How a Click-to-WhatsApp ad becomes a WhatsApp Pay transaction — without ever leaving the chat',
        13, False, C('C9D1D6'))], align=PP_ALIGN.CENTER)], align=PP_ALIGN.CENTER)
    stages = ['① AD', '② CHAT', '③ PAY', '④ ATTRIBUTION']
    cw, gap = 2.5, 0.25
    x0 = (13.333 - (4 * cw + 3 * gap)) / 2
    for i, st in enumerate(stages):
        pill(s, x0 + i * (cw + gap), 4.7, cw, 0.55, st, TEAL, WHITE, 13)
        if i < 3:
            tx(s, x0 + i * (cw + gap) + cw, 4.7, gap, 0.55, [P([R('→', 14, True, ACCENT)], align=PP_ALIGN.CENTER)],
               anchor=MSO_ANCHOR.MIDDLE)
    tx(s, 1.0, 5.75, 11.33, 0.5, [P([R(
        'Companion to the WhatsApp Pay × CTWA PRD  ·  Big Bet #1 / Tier-1 fix', 11, False, GREY)],
        align=PP_ALIGN.CENTER)], align=PP_ALIGN.CENTER)
    tx(s, 0.6, 6.85, 6.0, 0.4, [P([R('[Your Name] · PM Portfolio', 11, False, GREY)])], anchor=MSO_ANCHOR.MIDDLE)

    # -------- SLIDE 2 : THE CLOSED LOOP (overview) --------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    header(s, 'The Closed Loop — Why the Payment Matters', 'OVERVIEW'); footer(s, 2)
    # 3 boxes in a row + attribution box above
    boxes = [('① AD', 'Click-to-WhatsApp ad\non Instagram / Facebook', PRIMARY),
             ('② CHAT', 'In-chat catalog + cart\n(no app switch)', TEAL),
             ('③ PAY', 'WhatsApp Pay checkout\n(default) → UPI PIN', C('1AA890'))]
    bw, by = 3.4, 3.4
    x0 = 0.9
    centers = []
    for i, (t, d, col) in enumerate(boxes):
        bx = x0 + i * (bw + 0.55)
        centers.append(bx + bw / 2)
        rr(s, bx, by, bw, 1.25, col, radius=0.08)
        tx(s, bx + 0.15, by + 0.12, bw - 0.3, 1.0, [
            P([R(t, 15, True, WHITE)], align=PP_ALIGN.CENTER, sa=3),
            P([R(d, 11, False, C('EAFBF3'))], align=PP_ALIGN.CENTER, ln=1.05)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            glyph(s, bx + bw + 0.1, by + 0.45, 0.34, '→', WHITE, 15, fill=ACCENT)
    # attribution box on top spanning
    rr(s, centers[0] - 0.4, 1.65, centers[2] - centers[0] + 0.8, 0.95, DARK, radius=0.06)
    tx(s, centers[0] - 0.3, 1.72, centers[2] - centers[0] + 0.6, 0.85, [
        P([R('④ ATTRIBUTION', 13, True, ACCENT)], align=PP_ALIGN.CENTER, sa=2, ln=0.95),
        P([R('Each completed checkout fires a first-party purchase event → Ads optimization + advertiser ROAS',
             10.5, False, WHITE)], align=PP_ALIGN.CENTER, ln=1.05)],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrows: pay -> attribution -> ad (the loop)
    arrow(s, centers[2], by, centers[2], 2.6, ACCENT, 2.5)            # pay up to attribution
    arrow(s, centers[0], 2.6, centers[0], by, ACCENT, 2.5)            # attribution down to ad (loop back)
    tx(s, centers[2] + 0.05, 2.75, 1.8, 0.6, [P([R('payment = the\nconversion signal', 9, True, PRIMARY)], ln=1.0)])
    tx(s, centers[0] - 1.85, 2.75, 1.8, 0.6, [P([R('better signal →\nmore valuable ads', 9, True, PRIMARY)], align=PP_ALIGN.RIGHT, ln=1.0)])
    # bottom payoff
    rr(s, 0.7, 5.5, 11.93, 1.2, HILITE, line=ACCENT, lw=1.5)
    tx(s, 1.0, 5.6, 11.33, 1.0, [
        P([R('The point of the loop:  ', 13, True, PRIMARY), R(
            'under zero-MDR the payment earns ₹0 directly. Closing it inside WhatsApp turns each checkout into a measured ad-conversion — so the transaction finally has ',
            12.5, False, INK), R('positive LTV to Meta.', 12.5, True, INK)], ln=1.15),
        P([R('Template: WeChat monetizes a zero-MDR-style rail on ~0.6% commerce take over >RMB 2T GMV — not on the transaction.',
             10.5, False, SUBTLE, True)], sb=4)],
       anchor=MSO_ANCHOR.MIDDLE)

    # -------- SLIDE 3 : SCREENS 1-3 --------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    header(s, 'The Flow ①–③ : Ad → Chat → Cart', 'SCREENS'); footer(s, 3)
    cxs = [2.72, 6.66, 10.6]
    sb = phone(s, cxs[0] - 0.975, 1.3); draw_ad(s, *sb)
    sb = phone(s, cxs[1] - 0.975, 1.3); draw_chat_catalog(s, *sb)
    sb = phone(s, cxs[2] - 0.975, 1.3); draw_cart(s, *sb)
    pin(s, cxs[0] + 0.55, 4.15, 1); pin(s, cxs[1] + 0.5, 2.55, 2); pin(s, cxs[2] + 0.5, 3.72, 3)
    phone_caption(s, cxs[0], 1, 'The CTWA ad', "‘Shop on WhatsApp’ CTA. Ad & campaign ID ride the tap.")
    phone_caption(s, cxs[1], 2, 'Chat opens on the product', 'Native in-chat catalog — the conversation is the storefront.')
    phone_caption(s, cxs[2], 3, 'In-chat cart', 'Cart persists in the thread. One tap to checkout.')

    # -------- SLIDE 4 : SCREENS 4-6 --------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    header(s, 'The Flow ④–⑥ : Checkout → Pay → Confirm', 'SCREENS'); footer(s, 4)
    sb = phone(s, cxs[0] - 0.975, 1.3); draw_checkout(s, *sb)
    sb = phone(s, cxs[1] - 0.975, 1.3); draw_pin(s, *sb)
    sb = phone(s, cxs[2] - 0.975, 1.3); draw_success(s, *sb)
    pin(s, cxs[0] + 0.5, 3.0, 4, RED); pin(s, cxs[1] + 0.4, 2.55, 5); pin(s, cxs[2] + 0.5, 3.0, 6)
    phone_caption(s, cxs[0], 4, 'Checkout — Pay is the default', 'WhatsApp Pay pre-selected. Other UPI apps stay available.', RED)
    phone_caption(s, cxs[1], 5, 'UPI PIN, in context', 'NPCI PIN flow; returns straight to the chat. Amount is server-set.')
    phone_caption(s, cxs[2], 6, 'Receipt as a chat message', 'Order confirmation lives in the thread → a re-engagement surface.')

    # -------- SLIDE 5 : THE LOOP CLOSES (attribution) --------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    header(s, 'The Loop Closes: ₹0 Payment → Ad Attribution', 'THE PAYOFF'); footer(s, 5)
    # left: 3-step flow
    steps = [('✓  Payment completed in WhatsApp', ACCENT, WHITE, PRIMARY),
             ('→  First-party purchase event fires,\n     keyed to the CTWA ad / campaign', PRIMARY, WHITE, WHITE),
             ('→  Ads ranking optimizes + advertiser\n     sees ROAS & cost-per-purchase', TEAL, WHITE, WHITE)]
    sy0 = 1.5
    for i, (t, fill, _, tcol) in enumerate(steps):
        yy = sy0 + i * 1.25
        rr(s, 0.7, yy, 6.0, 1.0, fill, radius=0.08)
        tx(s, 0.95, yy, 5.6, 1.0, [P([R(t, 13, True, tcol)], ln=1.1)], anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            glyph(s, 3.5, yy + 1.0, 0.3, '↓', WHITE, 13, fill=DARK)
    # right: Ads Manager mock
    rr(s, 7.1, 1.5, 5.6, 3.5, WHITE, line=BORDER, radius=0.05)
    rc(s, 7.1, 1.5, 5.6, 0.5, DARK)
    tx(s, 7.3, 1.5, 5.2, 0.5, [P([R('◐  Meta Ads Manager · CTWA', 12, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
    rows = [('Campaign', 'Diwali Grocery'), ('Conversations', '1,240'),
            ('Purchases  (closed-loop)', '386'), ('Cost / purchase', '₹49'),
            ('Purchase value (GMV)', '₹2.3L'), ('ROAS', '6.2×')]
    ry = 2.15
    for i, (k, v) in enumerate(rows):
        if i % 2 == 1: rc(s, 7.2, ry, 5.4, 0.4, C('F7F9FA'))
        tx(s, 7.35, ry, 3.4, 0.4, [P([R(k, 11, False, INK)])], anchor=MSO_ANCHOR.MIDDLE)
        hot = k.startswith('Purchases') or k == 'ROAS'
        tx(s, 10.6, ry, 1.95, 0.4, [P([R(v, 12, True, (ACCENT if hot else INK) if not hot else PRIMARY)], align=PP_ALIGN.RIGHT)],
           anchor=MSO_ANCHOR.MIDDLE)
        ry += 0.42
    tx(s, 7.3, 4.62, 5.2, 0.3, [P([R('▲ 386 purchases now measured first-party — previously invisible', 9, True, ACCENT, True)])],
       anchor=MSO_ANCHOR.MIDDLE)
    # payoff banner
    rr(s, 0.7, 5.5, 11.93, 1.2, DARK, radius=0.06)
    tx(s, 1.0, 5.6, 11.33, 1.0, [
        P([R('This event is the entire point.  ', 14, True, ACCENT),
           R('A UPI payment worth ₹0 to Meta becomes a measured conversion that makes CTWA ads more effective — and more valuable.',
             12.5, False, WHITE)], ln=1.15),
        P([R('That is the positive LTV the diagnosis said was missing — without changing UPI economics at all.', 11, False, C('C9D1D6'), True)], sb=4)],
       anchor=MSO_ANCHOR.MIDDLE)

    # -------- SLIDE 6 : DESIGN PRINCIPLES --------
    s = prs.slides.add_slide(blank); bg(s, LIGHT)
    header(s, 'Key Design Decisions', 'PRINCIPLES'); footer(s, 6)
    cards = [
        ('◎', 'Default, not exclusive', 'WhatsApp Pay is pre-selected at checkout — but other UPI apps stay available, keeping the design defensible on self-preferencing.'),
        ('⇆', 'Never leave the chat', 'Ad → catalog → cart → pay → receipt all happen in one surface. Every app-switch is a drop-off removed.'),
        ('🔒', 'Server-set amounts', 'The payable amount comes from the cart, never user-typed — anti-fraud by construction.'),
        ('💬', 'Receipt = re-engagement', 'The order confirmation is a chat message the shopper returns to — a durable surface for tracking, support, and repeat purchase.'),
        ('📣', 'Acquisition = running an ad', 'Adopting the checkout is a side-effect of running a CTWA ad → merchant CAC ≈ ₹0, the structural edge over a field-force build.'),
        ('📈', 'Payment = measurement layer', 'Every checkout emits a first-party conversion event. The payment stops being a cost and becomes CTWA’s attribution engine.'),
    ]
    cw, ch, gx, gy = 3.85, 1.65, 0.3, 0.28
    x0, y0 = 0.7, 1.35
    cols = [PRIMARY, TEAL, C('1AA890'), PRIMARY, TEAL, C('1AA890')]
    for i, (g, t, d) in enumerate(cards):
        col = i % 3; row = i // 3
        bx = x0 + col * (cw + gx); byy = y0 + row * (ch + gy)
        rr(s, bx, byy, cw, ch, WHITE, line=BORDER, radius=0.06)
        rc(s, bx, byy, 0.12, ch, cols[i])
        glyph(s, bx + 0.22, byy + 0.16, 0.46, g, WHITE, 16, fill=cols[i])
        tx(s, bx + 0.8, byy + 0.14, cw - 0.95, ch - 0.25, [
            P([R(t, 13, True, INK)], sa=3, ln=1.0),
            P([R(d, 10, False, SUBTLE)], ln=1.08)])

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'WhatsApp_Pay_CTWA_Wireframes.pptx')
    prs.save(out)
    print('Saved:', out)
    print('Slides:', len(prs.slides._sldIdLst))
    for i, t in enumerate(['Title', 'The Closed Loop (overview)', 'Screens 1-3 (Ad/Chat/Cart)',
                           'Screens 4-6 (Checkout/PIN/Confirm)', 'The Loop Closes (attribution)',
                           'Key Design Decisions'], 1):
        print('  %d. %s' % (i, t))


if __name__ == '__main__':
    build()
